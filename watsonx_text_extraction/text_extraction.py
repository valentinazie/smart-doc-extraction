#!/usr/bin/env python3
"""
watsonx Text Extraction V2 — unified single-file / folder pipeline.
No LibreOffice dependency; watsonx V2 ingests .pdf/.docx/.pptx/.xlsx natively.

Usage:
    # Single file → results stay in COS
    python watsonx_text_extraction/text_extraction.py <file>

    # Folder (walked recursively) → results stay in COS
    python watsonx_text_extraction/text_extraction.py <folder>

    # ...plus local download when --output is given
    python watsonx_text_extraction/text_extraction.py <file|folder> --output ./out

    # Force re-extraction even if results already exist in COS
    python watsonx_text_extraction/text_extraction.py <file|folder> --reprocess

Examples (run from repo root, with the venv active):
    source venv/bin/activate

    python watsonx_text_extraction/text_extraction.py \\
        "client_data/data/F250904-5269_240828_VR8500배터리_연소.pdf"

    python watsonx_text_extraction/text_extraction.py \\
        client_data/data --output ./output/docx_results

Notes:
    - Reads credentials from ./.env via `common.config.load_env()`.
    - Extraction skips a file if a prior `assembly.md` exists in COS under the
      same prefix (`text_extraction_results/<file_stem>_`). Use --reprocess to
      force a re-run (old prefix is cleaned first).
    - Stem collision: two inputs that share a stem (e.g. `report.pdf` +
      `report.pptx`) collide on the same COS prefix and overwrite each other.
      Rename one (e.g. `report_pptx.pptx`) before running both.
    - Filenames with spaces are renamed in place (spaces -> underscores) before
      upload, so downloaded folder names also use underscores.
    - `.xls` / `.xlsm` are auto-converted to `.xlsx`; `.doc` / `.ppt` are NOT
      supported — convert to `.docx`/`.pptx` first.
    - Supported extensions: .pdf, .docx, .pptx, .xlsx, .xls, .xlsm
"""

import argparse
import os
import sys
import time
import json
from pathlib import Path

from ibm_watsonx_ai import APIClient, Credentials
from ibm_watsonx_ai.foundation_models.extractions import TextExtractionsV2, TextExtractionsV2ResultFormats
from ibm_watsonx_ai.helpers import DataConnection, S3Location
from ibm_watsonx_ai.metanames import TextExtractionsV2ParametersMetaNames

_SRC_ROOT = Path(__file__).resolve().parent.parent
if str(_SRC_ROOT) not in sys.path:
    sys.path.insert(0, str(_SRC_ROOT))

from common.config import (  # noqa: E402
    load_env,
    get_watsonx_credentials as get_credentials,
    get_space_cos_client,
)

load_env()

class NotebookTextExtraction:
    """Text extraction using pip-installable packages only"""
    
    def __init__(self):
        # Setup credentials and client (from notebook)
        self.space_id = os.environ["SPACE_ID"]
        self.credentials = get_credentials()
        self.watsonx_client = APIClient(credentials=self.credentials, space_id=self.space_id)
        try:
            self.cos_bucket_name = os.environ["COS_BUCKET_NAME"]
        except KeyError as exc:
            raise RuntimeError(
                "COS_BUCKET_NAME is required (set it in your .env, "
                "e.g. COS_BUCKET_NAME=my-watsonx-space-bucket)."
            ) from exc
        
        # Setup COS connection (from working analyzer)
        self.setup_cos()
        
        # Initialize extraction manager (from notebook)
        self.extraction = TextExtractionsV2(credentials=self.credentials, space_id=self.space_id)
    
    def setup_cos(self):
        """Setup COS client + connection for text extraction."""
        print("☁️  Setting up COS connection...")
        self.cos_client, _ = get_space_cos_client(
            self.watsonx_client, bucket=self.cos_bucket_name
        )

        buckets_names = [b["Name"] for b in self.cos_client.list_buckets()["Buckets"]]
        if self.cos_bucket_name not in buckets_names:
            self.cos_client.create_bucket(Bucket=self.cos_bucket_name)

        cos_credentials = self.watsonx_client.spaces.get_details(
            space_id=self.space_id)["entity"]["storage"]["properties"]
        connection_details = self.watsonx_client.connections.create({
            "datasource_type": self.watsonx_client.connections.get_datasource_type_uid_by_name("bluemixcloudobjectstorage"),
            "name": "Connection to COS for text extraction",
            "properties": {
                "bucket": self.cos_bucket_name,
                "access_key": cos_credentials["credentials"]["editor"]["access_key_id"],
                "secret_key": cos_credentials["credentials"]["editor"]["secret_access_key"],
                "iam_url": self.watsonx_client.service_instance._href_definitions.get_iam_token_url(),
                "url": cos_credentials["endpoint_url"],
            },
        })
        self.cos_connection_id = self.watsonx_client.connections.get_id(connection_details)
        print(f"✅ Connection ID: {self.cos_connection_id}")

    # Characters we strip from the filename stem before upload as defensive
    # hygiene against COS / shell path quirks. Spaces have always been here;
    # parens/brackets were added because they're commonly path-hostile in
    # storage layers, even though we never proved they caused a specific
    # watsonx failure (the failures we hit turned out to be a misread of an
    # intermediate "downloaded" status — see monitor_job).
    _UNSAFE_STEM_CHARS = " ()[]"

    @staticmethod
    def replace_spaces_in_filename(file_path: str) -> str:
        """Rename file in place to strip COS-unsafe characters from the stem.

        The real file extension (`.docx`, `.pdf`, …) is kept intact; only the
        stem portion is sanitized. Name kept for backward compatibility —
        historically only spaces were replaced. See _UNSAFE_STEM_CHARS for
        the full set we normalize.
        """
        directory, filename = os.path.split(file_path)
        if not os.path.exists(file_path):
            print(f"   ⚠️  File not found while renaming (already renamed?): {filename}")
            return file_path

        stem, ext = os.path.splitext(filename)
        translation = {ord(c): "_" for c in NotebookTextExtraction._UNSAFE_STEM_CHARS}
        new_stem = stem.translate(translation)
        # Collapse runs of underscores produced by adjacent unsafe chars (e.g. "()" → "__").
        while "__" in new_stem:
            new_stem = new_stem.replace("__", "_")
        new_stem = new_stem.strip("_")
        normalized = f"{new_stem}{ext}"
        if normalized == filename:
            return file_path

        candidate = os.path.join(directory, normalized)
        counter = 1
        while os.path.exists(candidate):
            candidate = os.path.join(directory, f"{new_stem}_{counter}{ext}")
            counter += 1

        os.rename(file_path, candidate)
        print(f"   ↪ Renamed '{filename}' → '{os.path.basename(candidate)}'")
        return candidate

    @staticmethod
    def convert_excel_to_xlsx(file_path: str) -> str:
        """Convert legacy/macro Excel formats (.xls/.xlsm) to .xlsx using pure Python libraries"""
        try:
            import pandas as pd
            import openpyxl
        except ImportError:
            raise RuntimeError(
                "Required packages not installed. Run: pip install pandas openpyxl xlrd"
            ) from None

        directory = os.path.dirname(file_path)
        base_name = os.path.splitext(os.path.basename(file_path))[0]
        output_path = os.path.join(directory, f"{base_name}.xlsx")

        if os.path.exists(output_path):
            print(f"   ⚠️  .xlsx version already exists for {base_name}, using it.")
            return output_path

        ext = os.path.splitext(file_path)[1].lower()
        print(f"   🔄 Converting {os.path.basename(file_path)} → .xlsx using pandas/openpyxl")

        try:
            if ext == ".xlsm":
                # XLSM can be read by openpyxl directly
                from openpyxl import load_workbook
                wb = load_workbook(file_path, keep_vba=False)
                wb.save(output_path)
            elif ext == ".xls":
                # XLS (old format) needs xlrd
                try:
                    df = pd.read_excel(file_path, engine='xlrd', sheet_name=None)
                    with pd.ExcelWriter(output_path, engine='openpyxl') as writer:
                        for sheet_name, sheet_df in df.items():
                            sheet_df.to_excel(writer, sheet_name=sheet_name, index=False)
                except ImportError:
                    raise RuntimeError("xlrd package required for .xls files. Run: pip install xlrd")
            else:
                raise RuntimeError(f"Unsupported Excel format: {ext}")

            if os.path.exists(output_path):
                backup_path = f"{file_path}.bak"
                os.rename(file_path, backup_path)
                print(f"   ✅ Created {os.path.basename(output_path)} (original saved as {os.path.basename(backup_path)})")
                return output_path
            else:
                raise RuntimeError(f"Conversion reported success but {output_path} was not created.")
        except Exception as e:
            raise RuntimeError(f"Excel conversion failed: {e}")

    @staticmethod
    def convert_office_to_pdf(file_path: str) -> str:
        """Convert DOCX/PPTX to PDF using pure Python libraries (text-only PDF)
        
        Note: .doc and .ppt (old formats) are not supported in pip-only version.
        Use the LibreOffice version or convert them to .docx/.pptx first.
        """
        ext = os.path.splitext(file_path)[1].lower()
        
        if ext in [".doc", ".ppt"]:
            raise RuntimeError(
                f"Old format {ext} is not supported in pip-only version. "
                "Please convert to .docx/.pptx first, or use the LibreOffice version."
            )
        
        directory = os.path.dirname(file_path)
        base_name = os.path.splitext(os.path.basename(file_path))[0]
        output_path = os.path.join(directory, f"{base_name}.pdf")

        if ext == ".docx":
            try:
                from docx import Document
                from reportlab.lib.pagesizes import letter
                from reportlab.pdfgen import canvas
            except ImportError:
                raise RuntimeError(
                    "Required packages not installed. Run: pip install python-docx reportlab"
                ) from None

            # Check if file exists
            if not os.path.exists(file_path):
                # Try absolute path
                abs_path = os.path.abspath(file_path)
                if os.path.exists(abs_path):
                    file_path = abs_path
                else:
                    raise RuntimeError(
                        f"File not found: {file_path}\n"
                        f"   Absolute path also not found: {abs_path}\n"
                        f"   Current working directory: {os.getcwd()}\n"
                        f"   Please verify the file path is correct."
                    )
            else:
                # Convert to absolute path to avoid path issues
                file_path = os.path.abspath(file_path)
            
            print(f"   🔄 Converting {os.path.basename(file_path)} → PDF (text-only) using python-docx + reportlab")
            print(f"   📁 Full path: {file_path}")
            
            try:
                doc = Document(file_path)
                c = canvas.Canvas(output_path, pagesize=letter)
                width, height = letter
                y = height - 50
                line_height = 14
                margin = 50

                for para in doc.paragraphs:
                    if para.text.strip():
                        text = para.text.strip()
                        # Simple text wrapping
                        words = text.split()
                        line = ""
                        for word in words:
                            test_line = line + (" " if line else "") + word
                            if c.stringWidth(test_line, "Helvetica", 10) > (width - 2 * margin):
                                if line:
                                    c.drawString(margin, y, line)
                                    y -= line_height
                                    if y < margin:
                                        c.showPage()
                                        y = height - margin
                                line = word
                            else:
                                line = test_line
                        if line:
                            c.drawString(margin, y, line)
                            y -= line_height * 1.5
                            if y < margin:
                                c.showPage()
                                y = height - margin

                c.save()
                print(f"   ✅ Created PDF: {os.path.basename(output_path)} (text-only, formatting may be lost)")
                return output_path
            except Exception as e:
                raise RuntimeError(f"DOCX to PDF conversion failed: {e}")

        elif ext in [".pptx", ".ppt"]:
            try:
                from pptx import Presentation
                from reportlab.lib.pagesizes import letter
                from reportlab.pdfgen import canvas
            except ImportError:
                raise RuntimeError(
                    "Required packages not installed. Run: pip install python-pptx reportlab"
                ) from None

            # Check if file exists
            if not os.path.exists(file_path):
                # Try absolute path
                abs_path = os.path.abspath(file_path)
                if os.path.exists(abs_path):
                    file_path = abs_path
                else:
                    raise RuntimeError(
                        f"File not found: {file_path}\n"
                        f"   Absolute path also not found: {abs_path}\n"
                        f"   Current working directory: {os.getcwd()}\n"
                        f"   Please verify the file path is correct."
                    )
            else:
                # Convert to absolute path to avoid path issues
                file_path = os.path.abspath(file_path)
            
            print(f"   🔄 Converting {os.path.basename(file_path)} → PDF (text-only) using python-pptx + reportlab")
            print(f"   📁 Full path: {file_path}")
            
            try:
                prs = Presentation(file_path)
                c = canvas.Canvas(output_path, pagesize=letter)
                width, height = letter
                margin = 50
                line_height = 14

                for slide_num, slide in enumerate(prs.slides, 1):
                    if slide_num > 1:
                        c.showPage()
                    y = height - margin

                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text.strip():
                            text = shape.text.strip()
                            words = text.split()
                            line = ""
                            for word in words:
                                test_line = line + (" " if line else "") + word
                                if c.stringWidth(test_line, "Helvetica", 10) > (width - 2 * margin):
                                    if line:
                                        c.drawString(margin, y, line)
                                        y -= line_height
                                        if y < margin:
                                            break
                                    line = word
                                else:
                                    line = test_line
                            if line and y >= margin:
                                c.drawString(margin, y, line)
                                y -= line_height * 1.5

                c.save()
                print(f"   ✅ Created PDF: {os.path.basename(output_path)} (text-only, formatting/images lost)")
                return output_path
            except Exception as e:
                raise RuntimeError(f"PPTX to PDF conversion failed: {e}")
        else:
            raise RuntimeError(f"Unsupported format for PDF conversion: {ext}")

    def has_existing_markdown_results(self, base_filename: str) -> bool:
        """Check COS for existing markdown outputs for this file."""
        return self.find_existing_prefix(base_filename) is not None

    def find_existing_prefix(self, base_filename: str) -> str | None:
        """Return the COS prefix of an existing extraction result, or None.

        The prefix has the form ``text_extraction_results/<stem>_<ts>/``.
        Picks the most-recently-modified match if several timestamps exist.
        """
        prefix = f"text_extraction_results/{base_filename}_"
        try:
            response = self.cos_client.list_objects_v2(
                Bucket=self.cos_bucket_name,
                Prefix=prefix,
                MaxKeys=200,
            )
        except Exception as exc:
            print(f"   ⚠️  Could not check existing results: {exc}")
            return None

        # Find the latest .md object and derive its parent prefix.
        md_objects = [obj for obj in response.get("Contents", [])
                      if obj["Key"].endswith(".md")]
        if not md_objects:
            return None
        md_objects.sort(key=lambda o: o.get("LastModified", 0), reverse=True)
        latest_key = md_objects[0]["Key"]
        # parent prefix = everything up to and including the last "/"
        parent = latest_key.rsplit("/", 1)[0] + "/"
        print(f"   📁 Found existing result: {latest_key}")
        return parent

    def create_parameters(self):
        """Create extraction parameters following notebook approach"""
        parameters = {
            TextExtractionsV2ParametersMetaNames.MODE: "high_quality",
            TextExtractionsV2ParametersMetaNames.OCR_MODE: "enabled",
            TextExtractionsV2ParametersMetaNames.LANGUAGES: ["en", "ko"],
            TextExtractionsV2ParametersMetaNames.AUTO_ROTATION_CORRECTION: True,
            TextExtractionsV2ParametersMetaNames.CREATE_EMBEDDED_IMAGES: "enabled_placeholder",
            TextExtractionsV2ParametersMetaNames.OUTPUT_DPI: 150,
            TextExtractionsV2ParametersMetaNames.OUTPUT_TOKENS_AND_BBOX: True,
        }
        return parameters

    def run_multiple_formats_extraction(self, input_filename, document_reference, results_reference, output_path_prefix):
        """Run multiple formats extraction following notebook approach"""
        print("\n🚀 RUNNING MULTIPLE FORMATS EXTRACTION")
        print("=" * 50)
        
        parameters = self.create_parameters()
        
        job_details = self.extraction.run_job(
            document_reference=document_reference,
            results_reference=results_reference,
            parameters=parameters,
            result_formats=[
                TextExtractionsV2ResultFormats.MARKDOWN,
            ],
        )
        
        job_id = self.extraction.get_job_id(extraction_details=job_details)
        print(f"📋 Job ID: {job_id}")
        
        job_details = self.monitor_job(job_id, "Multiple Formats")
        
        if job_details and job_details["entity"]["results"]["status"] == "completed":
            try:
                prefix_without_slash = output_path_prefix.rstrip('/')
                objects = None
                for attempt in range(12):
                    wait_time = 5
                    print(f"   ⏳ Waiting for results (attempt {attempt+1}/12, +{wait_time}s)...")
                    time.sleep(wait_time)
                    objects = self.cos_client.list_objects_v2(
                        Bucket=self.cos_bucket_name,
                        Prefix=prefix_without_slash
                    )
                    if 'Contents' in objects:
                        break
                
                if objects and 'Contents' in objects:
                    md_files = [obj['Key'] for obj in objects['Contents'] if obj['Key'].endswith(".md")]
                    if not md_files:
                        print("   ❌ No markdown files found in results.")
                        return False
                    print(f"   ✅ Markdown results available in COS under {prefix_without_slash}/")
                    for key in md_files:
                        print(f"      - {key}")
                    return True
                else:
                    print(f"   ❌ No files found with prefix: {prefix_without_slash}")
                    # Job reported success but wrote nothing under our prefix.
                    # Dump the relevant parts of job_details so we can see
                    # what watsonx actually thinks happened.
                    print("   🔎 Job details (results section):")
                    try:
                        results = job_details.get("entity", {}).get("results", {})
                        results_ref = job_details.get("entity", {}).get("results_reference", {})
                        print(f"      status:           {results.get('status')}")
                        print(f"      error:            {results.get('error')}")
                        print(f"      message:          {results.get('message')}")
                        print(f"      location:         {results.get('location')}")
                        print(f"      results_ref path: "
                              f"{results_ref.get('location', {}).get('path')}")
                    except Exception as e:
                        print(f"      (couldn't introspect job_details: {e})")
                    return False
            except Exception as e:
                print(f"❌ Error locating results: {e}")
                return False
        
        return False

    # watsonx Text Extraction reports several statuses; only `completed` and
    # `failed` are actually terminal. `downloaded` is an intermediate state
    # meaning "the service has downloaded the source file" — the job keeps
    # running after that. Treating `downloaded` as success caused us to bail
    # out early and miss the eventual `failed` status (with its error message).
    TERMINAL_STATES = {"completed", "failed", "canceled", "cancelled"}

    def monitor_job(self, job_id, job_name):
        """Monitor job until it reaches a terminal state."""
        print(f"⏳ Monitoring {job_name} job...")
        last_status = None

        while True:
            time.sleep(5)

            job_details = self.extraction.get_job_details(job_id)
            job_status = job_details["entity"]["results"]["status"]

            if job_status != last_status:
                print(f"\n[status: {job_status}]", end="", flush=True)
                last_status = job_status

            if job_status in self.TERMINAL_STATES:
                print()  # newline after the trailing dots
                if job_status == "completed":
                    return job_details
                print(f"❌ Job ended with status: {job_status}")
                err = job_details.get("entity", {}).get("results", {}).get("error")
                if err:
                    print(f"   error: {err}")
                return None

            print(".", end="", flush=True)


# ---------------------------------------------------------------------------
# Folder walker + per-file driver (merged from the old batch_extract.py)
# ---------------------------------------------------------------------------

SUPPORTED_EXTENSIONS = {".pdf", ".docx", ".pptx", ".xlsx", ".xlsm", ".xls"}
UNSUPPORTED_EXTENSIONS = {".doc", ".ppt"}


def collect_files(base_dir: str) -> list[str]:
    """Walk `base_dir` recursively and return absolute paths of supported docs."""
    files: list[str] = []
    for root, _dirs, filenames in os.walk(base_dir):
        for fname in sorted(filenames):
            if fname.startswith("."):
                continue
            ext = os.path.splitext(fname)[1].lower()
            if ext in UNSUPPORTED_EXTENSIONS:
                print(f"  ⏭️  Skipping unsupported legacy format: {fname}")
                continue
            if ext not in SUPPORTED_EXTENSIONS:
                continue
            files.append(os.path.join(root, fname))
    return files


def _fmt_duration(seconds: float) -> str:
    if seconds >= 3600:
        return f"{seconds / 3600:.1f} h"
    if seconds >= 60:
        return f"{seconds / 60:.1f} min"
    return f"{seconds:.1f} s"


def process_one_file(
    extractor: "NotebookTextExtraction",
    file_path: str,
    *,
    reprocess: bool = False,
    label: str = "",
) -> tuple[bool, str | None, str]:
    """Extract one file to COS.

    Returns ``(ok, cos_prefix, final_local_path)``:
        - ok: True if extraction succeeded or was skipped because results exist
        - cos_prefix: the `text_extraction_results/<stem>_<timestamp>/` prefix
          on COS (for downstream download), or None on failure/skip-reprocess-off
        - final_local_path: the on-disk path actually sent to COS (may differ
          from `file_path` if spaces were renamed or .xls/.xlsm was converted)
    """
    file_start = time.time()
    original_path = os.path.abspath(file_path)
    if not os.path.exists(original_path):
        print(f"❌ File not found: {original_path}")
        return False, None, original_path

    filename = os.path.basename(original_path)
    header = f"{label} {filename}" if label else filename
    print(f"\n{'=' * 70}")
    print(f"📄 {header}")
    print(f"{'=' * 70}")

    try:
        local_file_path = NotebookTextExtraction.replace_spaces_in_filename(original_path)
        if not os.path.exists(local_file_path):
            print(f"❌ File not found after space-normalization: {local_file_path}")
            return False, None, local_file_path

        ext = os.path.splitext(local_file_path)[1].lower()
        if ext in UNSUPPORTED_EXTENSIONS:
            print(f"❌ {ext} is not supported — convert to .docx/.pptx first.")
            return False, None, local_file_path

        if ext in {".xlsm", ".xls"}:
            converted = NotebookTextExtraction.convert_excel_to_xlsx(local_file_path)
            if not converted or not os.path.exists(converted):
                print(f"❌ Failed to convert legacy Excel to .xlsx")
                return False, None, local_file_path
            local_file_path = converted
            print(f"📄 Converted → {os.path.basename(local_file_path)}")

        filename = os.path.basename(local_file_path)
        base_filename = os.path.splitext(filename)[0]

        existing_prefix = extractor.find_existing_prefix(base_filename)
        if existing_prefix is not None:
            if reprocess:
                print(f"♻️  Results exist in COS for {base_filename} — --reprocess set, re-running.")
                try:
                    old = extractor.cos_client.list_objects_v2(
                        Bucket=extractor.cos_bucket_name,
                        Prefix=f"text_extraction_results/{base_filename}_",
                    )
                    if "Contents" in old:
                        for obj in old["Contents"]:
                            extractor.cos_client.delete_object(
                                Bucket=extractor.cos_bucket_name, Key=obj["Key"]
                            )
                        print(f"   🧹 Cleaned {len(old['Contents'])} stale result file(s)")
                except Exception as e:
                    print(f"   ⚠️  Cleanup warning: {e}")
            else:
                print(f"⚠️  Results already exist in COS for {base_filename}. "
                      f"Reusing them. (use --reprocess to force re-extraction)")
                # Return the existing prefix so callers with --output can still
                # download the cached result locally.
                return True, existing_prefix, local_file_path

        extractor.cos_client.upload_file(local_file_path, extractor.cos_bucket_name, filename)
        print(f"📤 Uploaded: {filename}")

        document_reference = DataConnection(
            connection_asset_id=extractor.cos_connection_id,
            location=S3Location(bucket=extractor.cos_bucket_name, path=filename),
        )
        document_reference.set_client(extractor.watsonx_client)

        timestamp = int(time.time())
        cos_prefix = f"text_extraction_results/{base_filename}_{timestamp}/"
        results_reference = DataConnection(
            connection_asset_id=extractor.cos_connection_id,
            location=S3Location(bucket=extractor.cos_bucket_name, path=cos_prefix),
        )
        results_reference.set_client(extractor.watsonx_client)
        print(f"📁 Output prefix: {cos_prefix}")

        success = extractor.run_multiple_formats_extraction(
            filename, document_reference, results_reference, cos_prefix
        )

        try:
            extractor.cos_client.delete_object(Bucket=extractor.cos_bucket_name, Key=filename)
        except Exception as e:
            print(f"⚠️  Could not remove uploaded source file from COS: {e}")

        if success:
            print(f"✅ Done: {filename}")
            return True, cos_prefix, local_file_path
        print(f"💥 Failed: {filename}")
        return False, None, local_file_path

    except Exception as e:
        print(f"❌ Error processing {filename}: {e}")
        import traceback
        traceback.print_exc()
        return False, None, original_path
    finally:
        print(f"⏱️  Duration: {_fmt_duration(time.time() - file_start)}")


def download_cos_prefix(
    extractor: "NotebookTextExtraction",
    cos_prefix: str,
    output_dir: str,
    *,
    local_subdir: str | None = None,
) -> int:
    """Download every object under `cos_prefix` into `output_dir/<local_subdir>/`.

    If `local_subdir` is None, the in-COS folder name (e.g. `<stem>_<ts>/`)
    is preserved on disk. Pass `local_subdir=stem` to drop the timestamp and
    align with the other pipelines (`<output>/<stem>/`).
    """
    bucket = extractor.cos_bucket_name
    cos_client = extractor.cos_client
    downloaded = 0
    continuation = None
    while True:
        kwargs = {"Bucket": bucket, "Prefix": cos_prefix, "MaxKeys": 1000}
        if continuation:
            kwargs["ContinuationToken"] = continuation
        resp = cos_client.list_objects_v2(**kwargs)
        for obj in resp.get("Contents", []):
            key = obj["Key"]
            if local_subdir is not None:
                relative = f"{local_subdir}/{key.removeprefix(cos_prefix)}"
            else:
                relative = key.removeprefix("text_extraction_results/")
            local = Path(output_dir) / relative
            local.parent.mkdir(parents=True, exist_ok=True)
            cos_client.download_file(bucket, key, str(local))
            downloaded += 1
        if resp.get("IsTruncated"):
            continuation = resp["NextContinuationToken"]
        else:
            break
    return downloaded


def main():
    overall_start = time.time()
    parser = argparse.ArgumentParser(
        description="watsonx Text Extraction V2 — accepts a file OR a folder (recursive).",
    )
    parser.add_argument(
        "path",
        help="Document file OR folder (walked recursively). "
             "Supported: .pdf, .docx, .pptx, .xlsx, .xls, .xlsm",
    )
    parser.add_argument(
        "--output", "-o",
        help="Local directory to download extracted results into after each "
             "successful extraction. If omitted, results stay in COS only.",
    )
    parser.add_argument(
        "--reprocess", action="store_true",
        help="Re-run extraction even if results already exist in COS for a file.",
    )
    args = parser.parse_args()

    print("🚀 WATSONX TEXT EXTRACTION V2")
    print("=" * 70)

    input_path = os.path.abspath(args.path)
    if not os.path.exists(input_path):
        print(f"❌ Path not found: {input_path}")
        return

    if os.path.isdir(input_path):
        files = collect_files(input_path)
        if not files:
            print(f"No supported documents found under: {input_path}")
            return
        print(f"Source directory: {input_path}")
        print(f"Found {len(files)} document(s):")
        for i, f in enumerate(files, 1):
            size_mb = os.path.getsize(f) / (1024 * 1024)
            print(f"  {i:3d}. {os.path.basename(f)}  ({size_mb:.1f} MB)")
    else:
        ext = os.path.splitext(input_path)[1].lower()
        if ext not in SUPPORTED_EXTENSIONS and ext not in UNSUPPORTED_EXTENSIONS:
            print(f"❌ Unsupported extension: {ext}")
            return
        files = [input_path]

    output_dir = None
    if args.output:
        output_dir = os.path.abspath(args.output)
        Path(output_dir).mkdir(parents=True, exist_ok=True)
        print(f"📥 Results will be downloaded to: {output_dir}")

    print("\n🔌 Initializing watsonx extractor...")
    extractor = NotebookTextExtraction()

    success_count = 0
    fail_count = 0
    skip_count = 0
    downloaded_files = 0

    total = len(files)
    for i, file_path in enumerate(files, 1):
        label = f"[{i}/{total}]"
        ok, cos_prefix, final_local_path = process_one_file(
            extractor, file_path, reprocess=args.reprocess, label=label
        )
        if ok:
            if cos_prefix is None:
                # process_one_file returns (True, None, ...) only for hard
                # failures of the existing-prefix lookup (rare). Treat as skip.
                skip_count += 1
            else:
                success_count += 1
                if output_dir:
                    try:
                        # Use the file stem so DOCX results land in
                        # `<output>/<stem>/` instead of the COS-side
                        # `<stem>_<timestamp>/` folder name. Keeps the layout
                        # consistent with PDF/PPTX/Excel pipelines.
                        stem = Path(final_local_path).stem
                        n = download_cos_prefix(
                            extractor, cos_prefix, output_dir,
                            local_subdir=stem,
                        )
                        downloaded_files += n
                        print(f"📥 Downloaded {n} file(s) → {output_dir}/{stem}/")
                    except Exception as e:
                        print(f"⚠️  Download failed for {cos_prefix}: {e}")
        else:
            fail_count += 1
        if total > 1:
            time.sleep(2)

    print("\n" + "=" * 70)
    print("SUMMARY")
    print("=" * 70)
    print(f"  ✅ Extracted: {success_count}")
    if skip_count:
        print(f"  ⏭️  Skipped (already in COS): {skip_count}")
    print(f"  ❌ Failed:    {fail_count}")
    if output_dir:
        print(f"  📥 Downloaded files: {downloaded_files} → {output_dir}")
    else:
        print(f"  📦 Results remain in COS under prefix: text_extraction_results/")
        print(f"     (run with --output DIR to download them locally)")
    print(f"  ⏱️  Total runtime: {_fmt_duration(time.time() - overall_start)}")
    print("=" * 70)


if __name__ == "__main__":
    main()

