"""
Spatial box extraction and shape analysis.
"""

import os
import json
from pathlib import Path
from datetime import datetime

try:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

from utils.geometry import (
    calculate_shapes_bounds,
    calculate_spatial_containment,
    get_group_center,
    simple_boxes_overlap,
)


class SpatialMixin:
    """Methods for spatial analysis of slide elements."""

    def extract_spatial_structure_native(self, pptx_path):
        """Extract spatial structure (positions) using native python-pptx"""
        try:
            prs = Presentation(pptx_path)
            
            spatial_structure = {
                'file_info': {
                    'name': Path(pptx_path).name,
                    'total_slides': len(prs.slides),
                    'slide_width': prs.slide_width,
                    'slide_height': prs.slide_height
                },
                'slides': []
            }
            
            if not prs.slides:
                print("   ❌ No slides found in PPTX")
                return None
            for slide_num, slide in enumerate(prs.slides, 1):
                slide_spatial = {
                    'slide_number': slide_num,
                    'shapes': []
                }
                
                for shape_idx, shape in enumerate(slide.shapes):
                    shape_info = {
                        'shape_index': shape_idx,
                        'position': {
                            'left': shape.left,
                            'top': shape.top,
                            'width': shape.width,
                            'height': shape.height
                        },
                        'shape_type': self.get_shape_type_name(shape.shape_type),
                        'has_text_frame': hasattr(shape, 'text_frame') and shape.text_frame is not None
                    }
                    slide_spatial['shapes'].append(shape_info)
                
                spatial_structure['slides'].append(slide_spatial)
            
            return spatial_structure
            
        except Exception as e:
            print(f"   ❌ Error extracting spatial structure: {e}")
            return None



    def extract_spatial_analysis(self, pptx_path, watsonx_content=None):
        """Extract spatial box information with positions"""
        print(f"\n📦 EXTRACTING SPATIAL ANALYSIS")
        print("=" * 50)
        
        if not PPTX_AVAILABLE:
            print("❌ python-pptx not available")
            return False
        
        try:
            prs = Presentation(pptx_path)
            
            # Extract metadata
            metadata = {
                'file_path': str(pptx_path),
                'file_type': 'pptx',
                'total_slides': len(prs.slides),
                'slide_width': prs.slide_width,
                'slide_height': prs.slide_height,
                'title': getattr(prs.core_properties, 'title', 'Unknown'),
                'author': getattr(prs.core_properties, 'author', 'Unknown')
            }
            
            self.comprehensive_data['metadata'] = metadata
            
            spatial_analysis = {
                'slides': [],
                'captured_shapes': []
            }
            
            if not prs.slides:
                print("   ❌ No slides found in PPTX")
                return False
            
            for slide_idx, slide in enumerate(prs.slides, 1):
                slide_watsonx = watsonx_content if slide_idx == 1 else None
                print(f"   📦 Processing slide {slide_idx}/{len(prs.slides)} spatial analysis...")
                slide_data = self.extract_slide_spatial_boxes(slide, slide_idx, slide_watsonx)
                spatial_analysis['slides'].append(slide_data)
                
                if self.config['generate_visual_captures']:
                    self.capture_target_shapes_both_types(slide, slide_idx, prs)
            
            self.comprehensive_data['spatial_analysis'] = spatial_analysis
            
            # Save spatial analysis
            spatial_path = self.spatial_analysis_dir / "spatial_analysis.json"
            with open(spatial_path, 'w', encoding='utf-8') as f:
                json.dump(spatial_analysis, f, ensure_ascii=False, indent=2)
            
            # Save local sectioning analysis if any lines were found
            self.save_local_sectioning_analysis(spatial_analysis['slides'])
            
            # Create local sectioning visualization (if enabled)
            if self.config['generate_local_sectioning_visualizations']:
                self.create_local_sectioning_visualization(spatial_analysis['slides'])
            
            # Create note about visual captures if none were made (only if visual captures are enabled)
            visual_captures_list = self.comprehensive_data.get('visual_captures', [])
            if self.config['generate_visual_captures'] and len(visual_captures_list) == 0:
                # Ensure visual_captures_dir exists
                self.visual_captures_dir.mkdir(parents=True, exist_ok=True)
                print(f"      ⚠️  DEBUG: No visual captures found in comprehensive_data")
                print(f"      ⚠️  DEBUG: slide_images status: {len(self.slide_images) if hasattr(self, 'slide_images') else 'No slide_images'}")
                visual_note_path = self.visual_captures_dir / "visual_captures_info.txt"
                with open(visual_note_path, 'w', encoding='utf-8') as f:
                    f.write("VISUAL CAPTURES STATUS\n")
                    f.write("=" * 30 + "\n\n")
                    f.write("❌ No visual captures were created.\n\n")
                    f.write("REASON:\n")
                    f.write("PPTX to PDF conversion failed or PDF processing failed.\n\n")
                    f.write("WHAT WAS ATTEMPTED:\n")
                    f.write("1. ✅ Looked for existing PDF file alongside PPTX\n")
                    f.write("2. ✅ Attempted automatic PPTX → PDF conversion using LibreOffice\n")
                    f.write("3. ❌ PDF processing or conversion failed\n\n")
                    f.write("TO TROUBLESHOOT:\n")
                    f.write("1. Check if LibreOffice is properly installed\n")
                    f.write("2. Manually export your PowerPoint as PDF and place it in the same directory\n")
                    f.write("3. Ensure PDF file has the same name as PPTX but with .pdf extension\n")
                    f.write("4. Re-run the analysis\n\n")
                    f.write("WHAT YOU'RE MISSING:\n")
                    f.write("• High-quality visual crops of tables, images, and groups\n")
                    f.write("• Enhanced visual mapping in summaries\n")
                    f.write("• PDF-based precise positioning for captures\n\n")
                    f.write("WHAT STILL WORKS:\n")
                    f.write("• Text structure extraction\n")
                    f.write("• Spatial analysis and positioning\n")
                    f.write("• Smart grouping and reading order\n")
                    f.write("• All analysis and summaries\n")
                print(f"      📝 Visual captures info saved: {visual_note_path}")
            
            print(f"✅ Spatial analysis extracted: {len(prs.slides)} slides")
            print(f"📸 Visual captures: {len(self.comprehensive_data.get('visual_captures', []))}")
            return True
            
        except Exception as e:
            print(f"❌ Error extracting spatial analysis: {e}")
            return False



    def extract_slide_spatial_boxes(self, slide, slide_number, watsonx_content=None):
        """Extract all elements using hierarchical group-based spatial mapping"""
        slide_data = {
            'slide_number': slide_number,
            'boxes': [],
            'total_shapes': len(slide.shapes),
            'spatial_map': {},
            'local_sections': [],
            'line_dividers': []
        }
        
        # Extract title if exists
        try:
            if hasattr(slide, 'shapes') and hasattr(slide.shapes, 'title'):
                title_shape = slide.shapes.title
                if title_shape and hasattr(title_shape, 'text'):
                    slide_data['title'] = title_shape.text
        except:
            slide_data['title'] = f'Slide {slide_number}'
        
        # STEP 1: Extract ALL elements with precise spatial mapping
        content_elements = []
        line_elements = []
        all_elements = []
        
        for shape_idx, shape in enumerate(slide.shapes):
            element_info = self.extract_shape_spatial_info(shape, shape_idx, watsonx_content)
            if element_info:
                all_elements.append(element_info)
                # Enhanced line detection: explicit lines + line-like shapes
                is_line = False
                # Method 1: Explicit line shapes
                if shape.shape_type == MSO_SHAPE_TYPE.LINE:
                    is_line = True
                    print(f"         📏 Found explicit LINE shape: {element_info['box_id']}")
                # Method 2: Line-like shapes based on aspect ratio and size
                else:
                    width = element_info['position']['width']
                    height = element_info['position']['height']
                    aspect_ratio = width / height if height > 0 else float('inf')
                    
                    # Detect thin vertical lines (very tall and narrow)
                    if width <= 15 and height >= 50 and aspect_ratio < 0.3:
                        is_line = True
                        print(f"         📏 Found line-like VERTICAL shape: {element_info['box_id']} ({width:.0f}×{height:.0f}, ratio={aspect_ratio:.2f})")
                    
                    # Detect thin horizontal lines (very wide and short)
                    elif height <= 15 and width >= 50 and aspect_ratio > 3.0:
                        is_line = True
                        print(f"         📏 Found line-like HORIZONTAL shape: {element_info['box_id']} ({width:.0f}×{height:.0f}, ratio={aspect_ratio:.2f})")
                # Classify as line or content
                if is_line:
                    line_elements.append(element_info)
                    slide_data['line_dividers'].append(element_info)
                else:
                    content_elements.append(element_info)
        
        print(f"   🗺️  Extracted {len(content_elements)} content + {len(line_elements)} lines")
        
        # DEBUG: Show ALL shapes to see if vertical line is there
        print(f"   🔍 DEBUG: All shapes detected:")
        all_shapes = content_elements + line_elements
        for shape in all_shapes:
            width = shape['position']['width']
            height = shape['position']['height'] 
            aspect_ratio = width / height if height > 0 else float('inf')
            shape_type = shape.get('shape_type', 'Unknown')
            print(f"      🔍 {shape['box_id']}: {width:.0f}×{height:.0f} (ratio={aspect_ratio:.2f}) - {shape_type}")
            
            # Check if this could be a line-like shape
            if aspect_ratio > 10.0 or aspect_ratio < 0.1:
                print(f"         🎯 POTENTIAL LINE: Very high/low aspect ratio!")
            if width < 10 and height > 50:
                print(f"         🎯 POTENTIAL VERTICAL LINE: Thin and tall!")
            if height < 10 and width > 50:
                print(f"         🎯 POTENTIAL HORIZONTAL LINE: Thin and wide!")
        
        # STEP 1.5: Consolidate table components (treat tables as single entities)
        consolidated_elements = self.consolidate_table_components(content_elements)
        print(f"   📋 Consolidated {len(content_elements)} → {len(consolidated_elements)} elements (table consolidation)")
        
        # STEP 1.6: UNIFIED CONSOLIDATION FIRST - Merge overlapping images + 80% overlapping text + lines
        # This must happen BEFORE table overlap detection so we have meaningful consolidated visual units
        consolidated_elements = self.create_unified_groups(consolidated_elements, line_elements)
        print(f"   🎯 Created unified groups (overlapping images + 80% text overlap + lines) - BEFORE table processing")
        
        # STEP 1.7: Associate overlapping shapes with table cells (after unified groups created)
        table_elements = [elem for elem in consolidated_elements if elem.get('shape_type') == 'Table']
        
        # Create SMART HYBRID element list: unified groups + standalone elements only
        # Get IDs of elements that are part of unified groups
        unified_group_member_ids = set()
        for elem in consolidated_elements:
            if elem.get('box_type') == 'unified_group' and 'member_ids' in elem:
                unified_group_member_ids.update(elem['member_ids'])
                print(f"      📦 UnifiedGroup {elem.get('box_id')} contains members: {elem['member_ids']}")
        
        # Smart hybrid list: consolidated elements + ONLY meaningful standalone elements
        hybrid_elements = list(consolidated_elements)  # Start with consolidated (includes unified groups)
        
        consolidated_ids = {ce.get('box_id') for ce in consolidated_elements}

        for elem in all_elements:
            elem_id = elem.get('box_id')
            # Skip if already represented by a unified group (in consolidated_elements)
            # or already a member of one.
            if elem_id in consolidated_ids or elem_id in unified_group_member_ids:
                continue
            shape_type = elem.get('shape_type')
            if shape_type == 'AutoShape':
                # Keep any AutoShape that carries text (titles, section headers, body bullets).
                # Empty AutoShapes (decorative rectangles, backgrounds) are dropped.
                if elem.get('text', '').strip():
                    hybrid_elements.append(elem)
                    print(f"      ➕ Adding standalone AutoShape: {elem_id} - '{elem.get('text','').strip()[:50]}'")
            elif shape_type in ('Picture', 'TextBox'):
                hybrid_elements.append(elem)
                print(f"      ➕ Adding standalone element: {elem_id} ({shape_type})")

        print(f"      📋 Excluded {len(unified_group_member_ids)} unified group members: {list(unified_group_member_ids)}")
        
        print(f"   🔄 Created hybrid element list: {len(consolidated_elements)} consolidated + {len(hybrid_elements) - len(consolidated_elements)} individual = {len(hybrid_elements)} total")
        
        # STEP 1.6: Capture visual regions for unified groups BEFORE table processing (re-enabled per user request)
        print(f"   🔍 DEBUG - Checking unified group capture: generate_visual_captures={self.config.get('generate_visual_captures', False)}")
        print(f"   🔍 DEBUG - Found {len(consolidated_elements)} consolidated elements")
        unified_groups = [e for e in consolidated_elements if e.get('box_type') == 'unified_group']
        print(f"   🔍 DEBUG - Found {len(unified_groups)} unified groups: {[g.get('box_id') for g in unified_groups]}")
        
        if self.config['generate_visual_captures']:
            for element in consolidated_elements:
                if element.get('box_type') == 'unified_group':
                    print(f"   📸 Attempting to capture unified group: {element.get('box_id')}")
                    self.capture_unified_group_visual(element, slide_number)
        else:
            print(f"   🚫 Visual captures disabled - skipping unified group captures")
        
        # STEP 1.7: Associate overlapping elements with table cells using PRECISE PowerPoint XML grid
        all_consumed_shape_ids = []
        for table_elem in table_elements:
            # Pass hybrid list (unified groups + individual standalone elements)
            consumed_ids = self.associate_overlapping_shapes_with_table_cells(table_elem, hybrid_elements, slide_number)
            all_consumed_shape_ids.extend(consumed_ids)
        
        # Remove consumed shapes from consolidated_elements
        if all_consumed_shape_ids:
            original_count = len(consolidated_elements)
            consolidated_elements = [elem for elem in consolidated_elements 
                                   if elem.get('box_id') not in all_consumed_shape_ids]
            removed_count = original_count - len(consolidated_elements)
            print(f"   🔗 Associated overlapping unified groups with {len(table_elements)} table(s), removed {removed_count} consumed groups: {all_consumed_shape_ids}")
        else:
            print(f"   🔗 Associated overlapping unified groups with {len(table_elements)} table(s), no groups consumed")
        
        # STEP 2.1: Cell-based table reclassification (NEW APPROACH - moved before smart grouping)
        print(f"   🔍 Checking tables for cell-based reclassification...")
        
        reclassified_count = 0
        for element in consolidated_elements:
            if element.get('shape_type') == 'Table':
                box_id = element.get('box_id', '')
                
                # Get table cell structure
                cell_contents = element.get('cell_contents', [])
                if not cell_contents:
                    continue
                    
                # Determine table dimensions by analyzing ALL cell coordinates (not just cells with content)
                max_row = 0
                max_col = 0
                cells_with_text = []
                
                for cell in cell_contents:
                    # Consider ALL cells, not just those with content
                    row = cell.get('row', 0)
                    col = cell.get('col', 0)
                    max_row = max(max_row, row)
                    max_col = max(max_col, col)
                    
                    if cell.get('text', '').strip():
                        cells_with_text.append((row, col, cell.get('text', '').strip()))
                
                table_rows = max_row + 1
                table_cols = max_col + 1
                
                print(f"      🔍 Table {box_id}: {table_rows}x{table_cols} ({len(cells_with_text)} cells with text)")
                
                # NEW RULE 1: 1x1 table with text → Convert to TextBox
                if table_rows == 1 and table_cols == 1 and cells_with_text:
                    print(f"      📝 Converting 1x1 table {box_id} to TextBox (treat as text)")
                    element['shape_type'] = 'TextBox'
                    element['box_type'] = 'text'
                    element['text'] = cells_with_text[0][2]  # Get the text from (0,0)
                    element['has_text'] = True
                    reclassified_count += 1
                
                # NEW RULE 2: 1x2 table → Split into two separate groups (regardless of text content)
                elif table_rows == 1 and table_cols == 2:
                    print(f"      ✂️  Splitting 1x2 table {box_id} into two separate groups")
                    
                    # Get table position for splitting
                    table_pos = element.get('position', {})
                    table_left = table_pos.get('left', 0)
                    table_top = table_pos.get('top', 0)
                    table_width = table_pos.get('width', 0)
                    table_height = table_pos.get('height', 0)
                    
                    # Try to use real_grid for precise cell boundaries
                    real_grid = element.get('real_grid')
                    
                    # Create two groups for each cell
                    created_groups = []
                    for col in range(2):  # (0,0) and (0,1)
                        if real_grid and real_grid.get('cells'):
                            # Use precise cell boundaries from PowerPoint XML
                            matching_cell = None
                            for cell_info in real_grid['cells']:
                                if cell_info.get('row') == 0 and cell_info.get('col') == col:
                                    matching_cell = cell_info
                                    break
                            
                            if matching_cell:
                                cell_position = {
                                    'left': matching_cell['left'],
                                    'top': matching_cell['top'],
                                    'width': matching_cell['width'],
                                    'height': matching_cell['height']
                                }
                                print(f"         📐 REAL cell({0},{col}): {matching_cell}")
                            else:
                                # Fallback to 50/50 split
                                cell_width = table_width // 2
                                cell_position = {
                                    'left': table_left + (col * cell_width),
                                    'top': table_top,
                                    'width': cell_width,
                                    'height': table_height
                                }
                                print(f"         📐 FALLBACK cell({0},{col}): {cell_position}")
                        else:
                            # Fallback to 50/50 split if no real_grid
                            cell_width = table_width // 2
                            cell_position = {
                                'left': table_left + (col * cell_width),
                                'top': table_top,
                                'width': cell_width,
                                'height': table_height
                            }
                            print(f"         📐 ESTIMATED cell({0},{col}): {cell_position}")
                        
                        # Find text for this cell
                        cell_text = ""
                        
                        for cell in cell_contents:
                            if cell.get('row') == 0 and cell.get('col') == col:
                                cell_text = cell.get('text', '').strip()
                                break
                        
                        if not cell_text:
                            cell_text = f"[Empty cell ({0},{col})]"
                        
                        # Create new UnifiedGroup for this cell - use UG prefix to avoid confusion with spatial groups
                        if not hasattr(self, '_next_group_id'):
                            self._next_group_id = 1  # Start from UG1
                        group_id = f"UG{self._next_group_id}"
                        self._next_group_id += 1
                        new_group = {
                            'box_id': group_id,
                            'shape_type': 'UnifiedGroup',
                            'box_type': 'unified_group',
                            'text': cell_text,
                            'has_text': bool(cell_text and cell_text != f"[Empty cell ({0},{col})]"),
                            'position': cell_position,
                            'is_standalone_group': (col == 0),  # Only left cell is standalone
                            'is_target_for_capture': True,
                            'visual_capture': f"slide_{slide_number:02d}_{group_id.lower()}_visual.png",
                            'original_table_id': box_id,
                            'cell_row': 0,
                            'cell_col': col
                        }
                        # Note: Don't pre-assign components - let spatial grouping logic find them based on boundaries
                        created_groups.append(new_group)
                        print(f"         ✅ Created {group_id}: '{cell_text}' (standalone={new_group['is_standalone_group']})")
                    
                    # Mark original table for removal
                    element['_mark_for_removal'] = True
                    element['_replacement_groups'] = created_groups
                    
                    reclassified_count += 1
                
                # All other cases remain as normal Table
                else:
                    print(f"      📋 Table {box_id} remains as normal Table ({table_rows}x{table_cols})")
                        
        if reclassified_count > 0:
            print(f"   🔄 Reclassified {reclassified_count} tables based on cell structure")
        else:
            print(f"   📋 No tables reclassified")
            
        # Process tables marked for removal and replacement
        tables_to_remove = []
        groups_to_add = []
        
        # Track converted table IDs for filtering from not-embedded lists
        if not hasattr(self, '_converted_table_ids'):
            self._converted_table_ids = []
            
        for element in consolidated_elements:
            if element.get('_mark_for_removal'):
                tables_to_remove.append(element)
                replacement_groups = element.get('_replacement_groups', [])
                groups_to_add.extend(replacement_groups)
                
                # Track this table ID as converted (for filtering from not-embedded lists)
                converted_table_id = element.get('box_id')
                if converted_table_id:
                    self._converted_table_ids.append(converted_table_id)
                    
                print(f"   ✂️  Replacing table {converted_table_id} with {len(replacement_groups)} groups: {[g.get('box_id') for g in replacement_groups]}")
        
        # Remove marked tables and add replacement groups
        for table in tables_to_remove:
            consolidated_elements.remove(table)
        
        for group in groups_to_add:
            consolidated_elements.append(group)
            
        if tables_to_remove:
            print(f"   🔄 Replaced {len(tables_to_remove)} tables with {len(groups_to_add)} groups")
            
        # STEP 2.2: Create smart groups AFTER table replacement (so it gets the new groups G2, G3, etc.)
        print(f"   🧠 Creating smart groups with updated elements (including replacement groups)...")
        smart_groups_result = self.create_smart_groups_for_slide(consolidated_elements, slide_number)
        
        # Handle case where no smart groups could be created
        if smart_groups_result is None or 'smart_groups' not in smart_groups_result:
            print(f"   ⚠️  No smart groups created, using individual components")
            # Create individual groups for each component
            smart_groups = []
            for idx, element in enumerate(content_elements):
                smart_groups.append({
                    'group_id': element['box_id'],
                    'members': [element],
                    'root_component': element,
                    'total_members': 1
                })
        else:
            # Convert smart_groups dictionary to list format expected by hierarchical code
            smart_groups_dict = smart_groups_result['smart_groups']
            smart_groups = []
            for group_id, group_data in smart_groups_dict.items():
                # Ensure root component has position structure
                root_component = group_data['root_component']
                if not isinstance(root_component, dict) or 'position' not in root_component:
                    print(f"      ⚠️  Skipping group {group_id}: invalid root component structure")
                    continue
                # Ensure all members have position structure
                valid_members = []
                for member_data in group_data['members']:
                    # Members have structure: {'index': X, 'box': {'position': {...}}}
                    if isinstance(member_data, dict) and 'box' in member_data:
                        member_box = member_data['box']
                        if isinstance(member_box, dict) and 'position' in member_box:
                            valid_members.append(member_box)  # Use the actual box data
                        else:
                            print(f"      ⚠️  Member box missing position in group {group_id}")
                    else:
                        print(f"      ⚠️  Invalid member structure in group {group_id}: {type(member_data)}")
                # Add root component to members list
                all_members = [root_component] + valid_members
                smart_groups.append({
                    'group_id': group_id,
                    'members': all_members,
                    'root_component': root_component,
                    'total_members': len(all_members)
                })
        
        print(f"   📦 Created {len(smart_groups)} smart groups for hierarchical reading")
        
        # STEP 3: Apply hierarchical group-based reading order
        if line_elements:
            # BYPASS grid-based sectional reading order - use direct TOP+LEFT reading order
            ordered_groups = self.sort_groups_by_simple_reading_order(smart_groups)
            slide_data['local_sections'] = self.create_hierarchical_sections(ordered_groups, line_elements)
        else:
            # No lines found, use simple group reading order
            print(f"   📖 No line dividers found, using simple group reading order")
            ordered_groups = self.sort_groups_by_simple_reading_order(smart_groups)
            slide_data['local_sections'] = [{
                'section_id': 'main',
                'section_type': 'full_slide_groups',
                'bounds': calculate_shapes_bounds(content_elements),
                'groups': ordered_groups,
                'reading_order': 1
            }]
        
        # STEP 4: Flatten groups to component list with hierarchical order
        slide_data['boxes'] = self.flatten_hierarchical_groups_to_components(ordered_groups, line_elements)
        
        # Store hierarchical groups for later use in visualization
        slide_data['hierarchical_groups'] = ordered_groups
        
        return slide_data



    def get_slide_bounds(self):
        """Get slide boundaries from metadata"""
        metadata = self.comprehensive_data.get('metadata', {})
        return {
            'top': 0,
            'left': 0,
            'right': metadata.get('slide_width', 9144000),  # Default PowerPoint width in EMU
            'bottom': metadata.get('slide_height', 6858000)  # Default PowerPoint height in EMU
        }



    def extract_shape_spatial_info(self, shape, shape_index, watsonx_content=None):
        """Extract comprehensive information from a PowerPoint shape"""
        try:
            # Basic shape information
            box_info = {
                'box_id': f"S{shape_index}",
                'shape_index': shape_index,
                'position': {
                    'left': shape.left,
                    'top': shape.top,
                    'width': shape.width,
                    'height': shape.height
                },
                'shape_type': self.get_shape_type_name(shape.shape_type),
                'text': '',
                'has_text': False,
                'is_target_for_capture': shape.shape_type in self.capture_target_types
            }
            
            # SPECIAL HANDLING FOR TABLES - Extract ALL content as single entity
            if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                table_content = self.extract_comprehensive_table_content(shape)
                # Check if this is a real table or just a layout container
                is_layout_container = (
                    table_content.get('table_rows', 0) <= 1 and 
                    table_content.get('table_cols', 0) <= 1 and 
                    not table_content.get('text', '').strip()
                )
                if is_layout_container:
                    # Treat 1x1 empty tables as regular shapes, not tables
                    box_info['box_type'] = 'layout_container'
                    box_info['shape_type'] = 'LayoutContainer'
                    box_info['layout_container_type'] = 'empty_table'
                    print(f"      📦 Detected layout container {box_info['box_id']}: {table_content.get('table_dimensions', 'unknown')} (empty table used as container)")
                else:
                    # PowerPoint tables are VISUAL structures - check actual table dimensions, not text delimiters
                    table_rows = table_content.get('table_rows', 0)
                    table_cols = table_content.get('table_cols', 0)
                    cell_contents = table_content.get('cell_contents', [])
                    
                    # Real table criteria: multiple rows OR columns OR multiple cells with content
                    is_real_table = (
                        table_rows > 1 or 
                        table_cols > 1 or 
                        len([cell for cell in cell_contents if cell.get('has_content', False)]) > 1
                    )
                    
                    if not is_real_table:
                        # Convert single-cell empty table to text box
                        table_text = table_content.get('text', '')
                        box_info['box_type'] = 'text'
                        box_info['shape_type'] = 'TextBox'
                        box_info['text'] = table_text
                        box_info['has_text'] = bool(table_text.strip())
                        box_info['is_target_for_capture'] = False
                        print(f"      📝 Converted single-cell table to TextBox {box_info['box_id']}: '{table_text[:50]}...' (no table structure)")
                    else:
                        # Real PowerPoint table with visual structure
                        box_info.update(table_content)
                        box_info['box_type'] = 'table'
                        box_info['shape_type'] = 'Table'
                        print(f"      📋 Keeping PowerPoint table {box_info['box_id']}: {table_rows}x{table_cols} with {len(cell_contents)} cells")
                        
                        # Extract real grid boundaries here (while we have the shape)
                        real_grid = self.extract_table_grid_boundaries(shape)
                        if real_grid:
                            box_info['real_grid'] = real_grid  # Store grid info (JSON serializable)
                            print(f"         📐 Extracted real table grid with {len(real_grid['cells'])} precise cells")
                        else:
                            print(f"         📐 Grid extraction failed, will use estimated boundaries")
                        
                        print(f"      📋 Extracted comprehensive table {box_info['box_id']}: {table_content.get('table_dimensions', 'unknown')} - {len(table_content.get('text', ''))} chars")
                return box_info
            
            # SPECIAL HANDLING FOR GROUPS - Extract group content
            elif shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                group_content = self.extract_comprehensive_group_content(shape)
                box_info.update(group_content)
                box_info['box_type'] = 'group'
                box_info['shape_type'] = 'Group'
                print(f"      📦 Extracted comprehensive group {box_info['box_id']}: {group_content.get('group_items_count', 0)} items - {len(group_content.get('text', ''))} chars")
                return box_info
            
            # HYBRID TEXT EXTRACTION: Try watsonx first, then PowerPoint fallback
            print(f"      🔍 ATTEMPTING WATSONX MATCH for shape with text: '{shape.text[:50] if hasattr(shape, 'text') and shape.text else 'NO_TEXT'}...'")
            watsonx_text = self.find_matching_watsonx_text(shape, watsonx_content)
            print(f"      📋 WATSONX MATCH RESULT: {'FOUND' if watsonx_text else 'NOT_FOUND'}")
            
            if watsonx_text:
                # Use high-quality watsonx text (preserves bullets, formatting)
                box_info['text'] = watsonx_text
                box_info['has_text'] = True
                box_info['text_source'] = 'watsonx'
                print(f"      ✅ Using watsonx text for {box_info['box_id']}: '{watsonx_text[:50]}...'")
            else:
                # Fallback to PowerPoint direct extraction
                if hasattr(shape, 'text'):
                    box_info['text'] = shape.text.strip()
                    box_info['has_text'] = bool(box_info['text'])
                elif hasattr(shape, 'text_frame') and shape.text_frame:
                    text_parts = []
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            text_parts.append(run.text)
                    box_info['text'] = ''.join(text_parts).strip()
                    box_info['has_text'] = bool(box_info['text'])
                box_info['text_source'] = 'powerpoint'
            
            # Determine box type
            box_info['box_type'] = self.get_box_type(shape.shape_type)
            
            return box_info
            
        except Exception as e:
            print(f"⚠️  Error extracting shape {shape_index}: {e}")
            return None



    def identify_vertical_lines(self, line_elements):
        """Identify vertical lines that can affect left-right ordering"""
        vertical_lines = []
        
        for line in line_elements:
            pos = line['position']
            width = max(pos['width'], 1)
            height = max(pos['height'], 1)
            aspect_ratio = width / height
            
            if aspect_ratio < 0.7:  # Vertical line
                line_center_x = pos['left'] + width // 2
                vertical_lines.append({
                    'line': line,
                    'center_x': line_center_x,
                    'top': pos['top'],
                    'bottom': pos['top'] + height
                })
                print(f"         📏 Vertical line {line['box_id']} at x={line_center_x}")
        
        # Sort vertical lines by X position
        vertical_lines.sort(key=lambda x: x['center_x'])
        return vertical_lines



    def find_primary_vertical_line_boundary(self, line_elements):
        """Find the most significant vertical line that acts as a content boundary"""
        vertical_lines = []
        
        for line in line_elements:
            pos = line['position']
            width = max(pos['width'], 1)
            height = max(pos['height'], 1)
            aspect_ratio = width / height
            
            # Only consider clear vertical lines as boundaries
            if aspect_ratio < 0.7:  # Vertical line
                line_score = height  # Longer lines are more significant boundaries
                vertical_lines.append((line, line_score))
                print(f"         📏 Vertical boundary candidate: {line['box_id']} (height={height}, score={line_score})")
        
        if not vertical_lines:
            return None
        
        # Return the longest (most significant) vertical line
        primary_line, best_score = max(vertical_lines, key=lambda x: x[1])
        print(f"         🎯 Selected primary vertical boundary: {primary_line['box_id']} (score={best_score})")
        
        return primary_line



    def detect_spatial_boundaries_in_groups(self, smart_groups):
        """Detect spatial boundaries within large groups (like tables with many components)"""
        spatial_boundaries = {'vertical': [], 'horizontal': []}
        
        for group in smart_groups:
            members = group.get('members', [])
            if len(members) < 8:  # Only analyze groups with many components
                continue
            
            group_id = group['group_id']
            print(f"         🔍 Analyzing spatial boundaries in large group {group_id} ({len(members)} components)")
            
            # Analyze X positions for vertical boundaries
            x_positions = []
            y_positions = []
            
            for member in members:
                if isinstance(member, dict) and 'box' in member:
                    pos = member['box']['position']
                else:
                    pos = member['position']
                x_center = pos['left'] + pos['width'] // 2
                y_center = pos['top'] + pos['height'] // 2
                x_positions.append(x_center)
                y_positions.append(y_center)
            
            # Find vertical boundaries (large X gaps)
            x_positions.sort()
            for i in range(len(x_positions) - 1):
                x_gap = x_positions[i + 1] - x_positions[i]
                if x_gap > 500000:  # Significant vertical gap
                    boundary_x = (x_positions[i] + x_positions[i + 1]) // 2
                    
                    # Create virtual vertical line
                    min_y = min(y_positions) - 100000
                    max_y = max(y_positions) + 100000
                    
                    virtual_line = {
                        'box_id': f'{group_id}_VBoundary_{len(spatial_boundaries["vertical"])+1}',
                        'position': {
                            'left': boundary_x,
                            'top': min_y,
                            'width': 0,
                            'height': max_y - min_y
                        },
                        'shape_type': 'VirtualLine',
                        'source': f'spatial_boundary_in_{group_id}'
                    }
                    spatial_boundaries['vertical'].append(virtual_line)
                    print(f"            📐 Virtual vertical boundary at X={boundary_x:,} (gap={x_gap:,}) in {group_id}")
            
            # Find horizontal boundaries (large Y gaps)
            y_positions.sort()
            for i in range(len(y_positions) - 1):
                y_gap = y_positions[i + 1] - y_positions[i]
                if y_gap > 300000:  # Significant horizontal gap
                    boundary_y = (y_positions[i] + y_positions[i + 1]) // 2
                    
                    # Create virtual horizontal line
                    min_x = min(x_positions) - 100000
                    max_x = max(x_positions) + 100000
                    
                    virtual_line = {
                        'box_id': f'{group_id}_HBoundary_{len(spatial_boundaries["horizontal"])+1}',
                        'position': {
                            'left': min_x,
                            'top': boundary_y,
                            'width': max_x - min_x,
                            'height': 0
                        },
                        'shape_type': 'VirtualLine',
                        'source': f'spatial_boundary_in_{group_id}'
                    }
                    spatial_boundaries['horizontal'].append(virtual_line)
                    print(f"            📐 Virtual horizontal boundary at Y={boundary_y:,} (gap={y_gap:,}) in {group_id}")
        
        total_virtual = len(spatial_boundaries['vertical']) + len(spatial_boundaries['horizontal'])
        if total_virtual > 0:
            print(f"         ✅ Detected {total_virtual} spatial boundaries ({len(spatial_boundaries['vertical'])} vertical, {len(spatial_boundaries['horizontal'])} horizontal)")
        
        return spatial_boundaries


