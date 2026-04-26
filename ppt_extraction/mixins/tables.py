"""
Table extraction, grid analysis, and consolidation.
"""

import os
import json
import time
from pathlib import Path
from datetime import datetime

try:
    from pptx.enum.shapes import MSO_SHAPE_TYPE
except ImportError:
    pass

from utils.geometry import (
    calculate_overlap_percentage,
    calculate_spatial_containment,
    calculate_spatial_overlap,
)


class TablesMixin:
    """Methods for extracting and processing table content."""

    def extract_tables_from_watsonx(self, watsonx_results):
        """Extract detailed table information from watsonx assembly results"""
        try:
            extracted_tables = []
            
            # Strategy 1: Look for explicit table structures
            if 'all_structures' in watsonx_results:
                all_structures = watsonx_results['all_structures']
                # Extract from tables structure if available
                if 'tables' in all_structures:
                    for table_idx, table in enumerate(all_structures['tables']):
                        table_data = {
                            'table_id': f'watsonx_table_{table_idx + 1}',
                            'type': 'structured_table',
                            'source': 'watsonx_assembly_tables',
                            'data': table
                        }
                        extracted_tables.append(table_data)
                # Strategy 2: Extract tabular content from tokens
                if 'tokens' in all_structures:
                    tokens = all_structures['tokens']
                    tabular_content = []
                    
                    # Look for patterns that suggest tabular data
                    for token in tokens:
                        if isinstance(token, dict) and 'text' in token:
                            text = token['text']
                            # Detect table patterns: "|", multiple spaces, numbers in columns
                            if ('|' in text or 
                                len(text.split()) > 3 or  # Multi-column data
                                any(char in text for char in [':', '=', '-']) or  # Table separators
                                (len(text.split()) >= 2 and any(part.replace('.','').replace(',','').isdigit() for part in text.split()[-2:]))):  # Numerical data
                                tabular_content.append(token)
                    
                    if tabular_content:
                        # Group consecutive tabular tokens into table structures
                        current_table = []
                        table_idx = 0
                        
                        for token in tabular_content:
                            current_table.append(token)
                            
                            # If we find a clear break or have enough content, save as table
                            if len(current_table) >= 3:  # Minimum table size
                                table_text = ' '.join([t.get('text', '') for t in current_table])
                                if len(table_text) > 20:  # Minimum meaningful content
                                    table_data = {
                                        'table_id': f'watsonx_token_table_{table_idx + 1}',
                                        'type': 'token_derived_table',
                                        'source': 'watsonx_assembly_tokens',
                                        'tokens': current_table.copy(),
                                        'reconstructed_text': table_text,
                                        'token_count': len(current_table)
                                    }
                                    extracted_tables.append(table_data)
                                    table_idx += 1
                                    current_table = []
                        
                        # Handle remaining tokens
                        if len(current_table) >= 2:
                            table_text = ' '.join([t.get('text', '') for t in current_table])
                            if len(table_text) > 15:
                                table_data = {
                                    'table_id': f'watsonx_token_table_{table_idx + 1}',
                                    'type': 'token_derived_table',
                                    'source': 'watsonx_assembly_tokens',
                                    'tokens': current_table,
                                    'reconstructed_text': table_text,
                                    'token_count': len(current_table)
                                }
                                extracted_tables.append(table_data)
            
            # Strategy 3: Check for Key-Value Pairs (KVPs) that might contain table data
            if 'kvps' in watsonx_results:
                kvp_tables = []
                for kvp in watsonx_results['kvps']:
                    kvp_str = str(kvp).lower()
                    if any(indicator in kvp_str for indicator in ['table', 'row', 'column', 'cell', 'data']):
                        kvp_tables.append(kvp)
                if kvp_tables:
                    table_data = {
                        'table_id': 'watsonx_kvp_table',
                        'type': 'kvp_derived_table',
                        'source': 'watsonx_assembly_kvps',
                        'data': kvp_tables,
                        'kvp_count': len(kvp_tables)
                    }
                    extracted_tables.append(table_data)
            
            # Strategy 4: Full content analysis for tabular patterns
            full_content = str(watsonx_results)
            if any(pattern in full_content.lower() for pattern in ['analysis item', 'freq range', '|', 'target', 'audit']):
                # This looks like it contains tabular data based on content
                table_data = {
                    'table_id': 'watsonx_content_analysis',
                    'type': 'content_derived_table',
                    'source': 'watsonx_full_content_analysis',
                    'content_indicators': [pattern for pattern in ['analysis item', 'freq range', 'target', 'audit'] if pattern in full_content.lower()],
                    'estimated_tabular_content': True
                }
                extracted_tables.append(table_data)
            
            print(f"   📊 Extracted {len(extracted_tables)} table structures from watsonx")
            if extracted_tables:
                for table in extracted_tables:
                    print(f"      📋 {table['table_id']}: {table['type']}")
            
            return extracted_tables
            
        except Exception as e:
            print(f"   ⚠️  Error extracting tables from watsonx: {e}")
            return []



    def save_table_extractions(self, table_data, slide_number):
        """Save table extraction data to files"""
        try:
            if not table_data:
                print(f"   📊 No tables found for slide {slide_number}")
                return
            
            # Save individual table files
            table_files = []
            for table in table_data:
                table_filename = f"slide_{slide_number:02d}_{table['table_id']}.json"
                table_path = self.table_extractions_dir / table_filename
                with open(table_path, 'w', encoding='utf-8') as f:
                    json.dump({
                        'slide_number': slide_number,
                        'extraction_timestamp': datetime.now().isoformat(),
                        'table_info': table
                    }, f, ensure_ascii=False, indent=2)
                table_files.append(table_filename)
                print(f"   💾 Saved table: {table_filename}")
            
            # Save summary file for the slide
            summary_filename = f"slide_{slide_number:02d}_tables_summary.json"
            summary_path = self.table_extractions_dir / summary_filename
            
            summary_data = {
                'slide_number': slide_number,
                'total_tables': len(table_data),
                'table_files': table_files,
                'table_types': [t['type'] for t in table_data],
                'extraction_method': 'watsonx_assembly_analysis'
            }
            
            with open(summary_path, 'w', encoding='utf-8') as f:
                json.dump(summary_data, f, ensure_ascii=False, indent=2)
            
            print(f"   📊 Table extraction summary: {len(table_data)} tables saved for slide {slide_number}")
            return table_files
            
        except Exception as e:
            print(f"   ❌ Error saving table extractions: {e}")
            return []



    def load_watsonx_table_data(self, slide_number):
        """Load real table boundaries from watsonx extraction files"""
        try:
            table_extractions_dir = self.output_dir / f'table_extractions'
            watsonx_tables = {}
            
            # Load ALL watsonx tables (structured + token tables)
            import glob
            
            # Load structured tables
            for table_num in [1, 2]:
                table_file = table_extractions_dir / f'slide_{slide_number:02d}_watsonx_table_{table_num}.json'
                if table_file.exists():
                    with open(table_file, 'r', encoding='utf-8') as f:
                        table_data = json.load(f)
                        table_info = table_data.get('table_info', {}).get('data', {})
                        if 'bbox_list' in table_info and table_info['bbox_list']:
                            bbox = table_info['bbox_list'][0]
                            watsonx_tables[f'STRUCT_TABLE_{table_num}'] = {
                                'bbox': bbox,
                                'rows': table_info.get('children_ids', []),
                                'table_id': table_info.get('id'),
                                'type': 'structured'
                            }
                            print(f'      📊 Loaded structured table {table_num}: {bbox}')
            
            # Load token tables (may have more table boundaries)
            token_files = list(table_extractions_dir.glob(f'slide_{slide_number:02d}_watsonx_token_table_*.json'))
            for token_file in token_files[:5]:  # Limit to first 5 to avoid overload
                try:
                    with open(token_file, 'r', encoding='utf-8') as f:
                        token_data = json.load(f)
                        tokens = token_data.get('table_info', {}).get('tokens', [])
                        if tokens:
                            # Calculate bounding box from token positions
                            x_coords = [t['bbox']['x'] for t in tokens if 'bbox' in t]
                            y_coords = [t['bbox']['y'] for t in tokens if 'bbox' in t]
                            if x_coords and y_coords:
                                min_x, max_x = min(x_coords), max(x_coords)
                                min_y, max_y = min(y_coords), max(y_coords)
                                token_bbox = {
                                    'x': min_x, 'y': min_y,
                                    'width': max_x - min_x + 100,  # Add some padding
                                    'height': max_y - min_y + 100,
                                    'page_number': tokens[0]['bbox']['page_number']
                                }
                                table_id = token_data.get('table_info', {}).get('table_id', token_file.stem)
                                watsonx_tables[table_id] = {
                                    'bbox': token_bbox,
                                    'tokens': tokens,
                                    'type': 'token_derived'
                                }
                                print(f'      📊 Loaded token table {table_id}: {token_bbox}')
                except Exception as e:
                    print(f'      ⚠️  Error loading {token_file.name}: {e}')
            
            return watsonx_tables
        except Exception as e:
            print(f'      ⚠️  Error loading watsonx table data: {e}')
            return {}



    def associate_overlapping_shapes_with_table_cells(self, table_element, all_elements, slide_number):
        """Associate spatially overlapping images/charts with specific table cells"""
        consumed_shape_ids = []
        
        try:
            if not table_element.get('cell_contents'):
                return consumed_shape_ids
                
            table_pos = table_element.get('position', {})
            table_rows = table_element.get('table_rows', 0)
            table_cols = table_element.get('table_cols', 0)
            
            if table_rows == 0 or table_cols == 0:
                return consumed_shape_ids
                
            # Load real watsonx table data for precise boundaries
            if not hasattr(self, 'watsonx_tables'):
                self.watsonx_tables = self.load_watsonx_table_data(slide_number)
            
            # Try to find corresponding watsonx table data
            table_id = table_element.get('box_id')
            watsonx_match = None
            
            # Match PowerPoint table with watsonx table based on position/size similarity  
            table_left = table_pos.get('left', 0)
            table_top = table_pos.get('top', 0)
            table_width = table_pos.get('width', 0)
            table_height = table_pos.get('height', 0)
            
            print(f'      🔍 PowerPoint table {table_id}: pos=({table_left}, {table_top}), size=({table_width}, {table_height})')
            
            best_match = None
            best_score = 0
            
            for watsonx_id, watsonx_data in self.watsonx_tables.items():
                watsonx_bbox = watsonx_data['bbox']
                watsonx_left = watsonx_bbox['x']
                watsonx_top = watsonx_bbox['y'] 
                watsonx_width = watsonx_bbox['width']
                watsonx_height = watsonx_bbox['height']
                
                print(f'         💡 Comparing with {watsonx_id}: pos=({watsonx_left}, {watsonx_top}), size=({watsonx_width}, {watsonx_height})')
                
                # Calculate similarity scores (more flexible matching)
                # Size similarity (most reliable indicator)
                size_sim_w = 1 - abs(table_width - watsonx_width) / max(table_width, watsonx_width, 1)
                size_sim_h = 1 - abs(table_height - watsonx_height) / max(table_height, watsonx_height, 1)
                size_score = (size_sim_w + size_sim_h) / 2
                
                # Relative position similarity (less strict for coordinate system differences)
                if table_width > 0 and table_height > 0:
                    rel_pos_x = abs(table_left / table_width - watsonx_left / watsonx_width) if watsonx_width > 0 else 1
                    rel_pos_y = abs(table_top / table_height - watsonx_top / watsonx_height) if watsonx_height > 0 else 1
                    pos_score = max(0, 1 - (rel_pos_x + rel_pos_y) / 2)
                else:
                    pos_score = 0
                
                # Combined score (prioritize size match)
                overall_score = size_score * 0.8 + pos_score * 0.2
                
                print(f'         📊 Scores: size={size_score:.2f}, pos={pos_score:.2f}, overall={overall_score:.2f}')
                
                if overall_score > best_score and overall_score > 0.3:  # Minimum threshold
                    best_match = watsonx_data
                    best_score = overall_score
                    print(f'         ✅ New best match: {watsonx_id} (score={best_score:.2f})')
            
            if best_match:
                watsonx_match = best_match
                print(f'      🎯 Selected best watsonx match for {table_id} with score {best_score:.2f}')
            else:
                print(f'      ❌ No suitable watsonx match found for {table_id} (best score: {best_score:.2f})')
            
            if watsonx_match:
                # Use watsonx table structure for precise cell mapping
                watsonx_bbox = watsonx_match['bbox']
                table_left = watsonx_bbox['x']  # Use watsonx coordinates
                table_top = watsonx_bbox['y']
                table_width = watsonx_bbox['width'] 
                table_height = watsonx_bbox['height']
                print(f'      ✅ Using watsonx precise boundaries: {watsonx_bbox}')
            else:
                # Fallback to estimated approach
                print(f'      ⚠️  No watsonx match found for {table_id}, using estimated cell boundaries')
            
            # STEP 1: Check if we have pre-extracted real table grid boundaries
            real_grid = table_element.get('real_grid')
            use_real_grid = real_grid is not None

            if use_real_grid:
                print(f"      ✅ Using REAL table grid with {len(real_grid['cells'])} precise cell boundaries")
                table_left = real_grid['table_bounds']['left']
                table_top = real_grid['table_bounds']['top']
                table_width = real_grid['table_bounds']['width']
                table_height = real_grid['table_bounds']['height']
                # Define fallback variables even when using real grid (for error cases)
                cell_width = table_width / table_cols if table_cols > 0 else 0
                cell_height = table_height / table_rows if table_rows > 0 else 0
            else:
                # Fallback to estimated approach if no real grid available
                print(f"      ⚠️  Using ESTIMATED cell boundaries (no real grid available)")
                # Calculate cell boundaries (fallback estimated approach)
                cell_width = table_width / table_cols if table_cols > 0 else 0
                cell_height = table_height / table_rows if table_rows > 0 else 0

            # Build a phantom-cell redirect map. python-pptx exposes merged cells as
            # one origin cell (with row_span/col_span > 1) PLUS the spanned-over cells
            # as separate phantom entries with their own bbox. The phantom bbox often
            # captures the lower half of the merged region, so a shape that visually
            # sits in the next data row gets misassigned to the phantom (the same
            # column's "row 1" sub-header strip, etc.). Redirect any phantom hit to
            # the closest non-phantom cell with content — preferring the next data
            # cell DOWN in the same column for row_span phantoms (so images in a
            # data row don't get pulled up into a merged header), then falling back
            # to the merged origin.
            phantom_to_redirect = {}
            if use_real_grid:
                phantom_cells = set()
                origin_of = {}  # phantom cell -> origin (row, col)
                for cell_info in real_grid['cells']:
                    rs = cell_info.get('row_span', 1) or 1
                    cs = cell_info.get('col_span', 1) or 1
                    if rs > 1 or cs > 1:
                        for dr in range(rs):
                            for dc in range(cs):
                                if dr == 0 and dc == 0:
                                    continue
                                key = (cell_info['row'] + dr, cell_info['col'] + dc)
                                phantom_cells.add(key)
                                origin_of[key] = (cell_info['row'], cell_info['col'])
                cell_lookup = {(c['row'], c['col']): c for c in (table_element.get('cell_contents') or [])}
                for (pr, pc) in phantom_cells:
                    target = None
                    # Prefer the next non-phantom cell with content DOWNWARD in the same column.
                    for tr in range(pr + 1, table_rows):
                        if (tr, pc) in phantom_cells:
                            continue
                        tc = cell_lookup.get((tr, pc))
                        if tc and tc.get('has_content') and (tc.get('text') or '').strip():
                            target = (tr, pc)
                            break
                    # Fallback to the merged origin (header cell).
                    if target is None:
                        target = origin_of.get((pr, pc))
                    if target and target != (pr, pc):
                        phantom_to_redirect[(pr, pc)] = target
                if phantom_to_redirect:
                    print(f"      🔁 Phantom-cell redirects: {phantom_to_redirect}")
            
            # Find ALL overlapping elements (images, charts, text boxes, autoshapes, graphics, inner tables)
            # ALLOW embedded tables but exclude the table itself to avoid self-referencing
            overlappable_elements = [elem for elem in all_elements 
                                   if elem.get('shape_type') not in ['BoundaryTable'] 
                                   and elem.get('box_id') != table_element.get('box_id')]
            
            # CRITICAL: Add the table's own contained_components for cell assignment
            if 'contained_components' in table_element:
                for contained_comp in table_element['contained_components']:
                    # Only add if not already in the list
                    if not any(elem.get('box_id') == contained_comp.get('box_id') for elem in overlappable_elements):
                        overlappable_elements.append(contained_comp)
                        print(f"         ➕ Added contained component {contained_comp.get('box_id')} for cell assignment")
            
            print(f"      🔍 Checking {len(overlappable_elements)} elements for overlap with table {table_element.get('box_id')}")
            for elem in overlappable_elements:
                print(f"         • {elem.get('box_id')} ({elem.get('shape_type')})")
            
            # ENHANCED CELL-BASED COMPONENT ASSIGNMENT
            # For each overlapping component, calculate its overlap with each cell and assign to best cell
            consumed_shape_ids = []
            not_embedded_components = []  # Components that overlap with table but not assigned to specific cells
            
            for component in overlappable_elements:
                comp_pos = component.get('position', {})
                comp_left = comp_pos.get('left', 0)
                comp_top = comp_pos.get('top', 0)
                comp_width = comp_pos.get('width', 0)
                comp_height = comp_pos.get('height', 0)
                
                print(f"      🔍 Analyzing {component.get('box_id')} for cell assignment")
                
                best_overlap = 0
                best_cell_row = -1
                best_cell_col = -1
                
                # Calculate overlap with each cell
                if use_real_grid:
                    # Use precise cell boundaries from real_grid
                    for cell_info in real_grid['cells']:
                        cell_row = cell_info['row']
                        cell_col = cell_info['col']
                        cell_left = cell_info['left']
                        cell_top = cell_info['top']
                        cell_right = cell_info['right']
                        cell_bottom = cell_info['bottom']
                        
                        # Calculate intersection
                        intersect_left = max(comp_left, cell_left)
                        intersect_top = max(comp_top, cell_top)
                        intersect_right = min(comp_left + comp_width, cell_right)
                        intersect_bottom = min(comp_top + comp_height, cell_bottom)
                        
                        if intersect_left < intersect_right and intersect_top < intersect_bottom:
                            intersect_area = (intersect_right - intersect_left) * (intersect_bottom - intersect_top)
                            comp_area = comp_width * comp_height
                            overlap_percentage = (intersect_area / comp_area) if comp_area > 0 else 0
                            
                            print(f"         📐 Cell({cell_row},{cell_col}): {overlap_percentage:.1%} overlap")
                        else:
                            overlap_percentage = 0
                        
                        if overlap_percentage > best_overlap:
                                best_overlap = overlap_percentage
                                best_cell_row = cell_row
                                best_cell_col = cell_col
                else:
                    # Use estimated cell boundaries
                    for row in range(table_rows):
                        for col in range(table_cols):
                            cell_left = table_left + (col * cell_width)
                            cell_top = table_top + (row * cell_height)
                            cell_right = cell_left + cell_width
                            cell_bottom = cell_top + cell_height
                            
                            # Calculate intersection
                            intersect_left = max(comp_left, cell_left)
                            intersect_top = max(comp_top, cell_top)
                            intersect_right = min(comp_left + comp_width, cell_right)
                            intersect_bottom = min(comp_top + comp_height, cell_bottom)
                            
                            if intersect_left < intersect_right and intersect_top < intersect_bottom:
                                intersect_area = (intersect_right - intersect_left) * (intersect_bottom - intersect_top)
                                comp_area = comp_width * comp_height
                                overlap_percentage = (intersect_area / comp_area) if comp_area > 0 else 0
                                
                                print(f"         📐 Cell({row},{col}): {overlap_percentage:.1%} overlap")
                                
                                if overlap_percentage > best_overlap:
                                    best_overlap = overlap_percentage
                                    best_cell_row = row
                                    best_cell_col = col
                
                # If the chosen cell is a merged-cell phantom, redirect the assignment.
                redirect = phantom_to_redirect.get((best_cell_row, best_cell_col))
                if redirect:
                    print(f"         🔁 Redirecting {component.get('box_id')} from phantom cell({best_cell_row},{best_cell_col}) → cell{redirect}")
                    best_cell_row, best_cell_col = redirect

                # Banner guard: an AutoShape (or TextBox) carrying text whose top sits
                # significantly above the table's top is almost certainly a header /
                # subtitle banner that is rendered OUTSIDE the table, even though its
                # bounding box dips into the top row. Don't let it be consumed by a
                # cell — it must remain a standalone group so its text is preserved.
                if (component.get('shape_type') in ('AutoShape', 'TextBox')
                        and component.get('text', '').strip()
                        and comp_height > 0):
                    above_table = max(0, table_top - comp_top)
                    above_frac = above_table / comp_height
                    if above_frac >= 0.20:
                        print(f"         🚫 Skipping {component.get('shape_type')} {component.get('box_id')} as banner ({above_frac*100:.1f}% of its height sits above table top)")
                        continue

                # Assign component to best cell - use higher threshold for TextBox to keep title texts as standalone groups
                overlap_threshold = 0.30 if component.get('shape_type') == 'TextBox' else 0.5
                if best_overlap > overlap_threshold:
                    print(f"         ✅ Assigning {component.get('box_id')} to cell({best_cell_row},{best_cell_col}) with {best_overlap:.1%} overlap")
                    
                    # Find the corresponding cell in cell_contents and add the component
                    for cell in table_element['cell_contents']:
                        if cell['row'] == best_cell_row and cell['col'] == best_cell_col:
                            if 'shapes' not in cell:
                                cell['shapes'] = []
                            
                            # Add component info to cell
                            # Generate visual capture path dynamically for Groups (both UG and G prefixes)
                            visual_capture_path = component.get('visual_capture', '')
                            group_id = component.get('box_id', '')
                            
                            if not visual_capture_path and group_id.startswith('UG'):
                                # New UnifiedGroup (UG1, UG2, UG3) - generate UG visual capture
                                try:
                                    group_index = int(group_id[2:])  # UG1 -> 1, UG2 -> 2
                                    visual_capture_path = f"slide_{slide_number:02d}_ug{group_index}_visual.png"
                                    print(f"         📸 Generated visual capture for {group_id}: {visual_capture_path}")
                                except (ValueError, IndexError):
                                    visual_capture_path = ''
                            elif not visual_capture_path and group_id.startswith('G') and group_id[1:].isdigit():
                                # Old spatial group (G1, G2, G3) - generate G visual capture
                                try:
                                    group_index = int(group_id[1:])  # G1 -> 1, G2 -> 2
                                    visual_capture_path = f"slide_{slide_number:02d}_g{group_index}_visual.png"
                                    print(f"         📸 Generated visual capture for {group_id}: {visual_capture_path}")
                                except (ValueError, IndexError):
                                    visual_capture_path = ''
                            
                            component_info = {
                                'box_id': component.get('box_id'),
                                'type': component.get('shape_type'),  # Use 'type' key to match expected format
                                'text': component.get('text', ''),
                                'overlap_percentage': best_overlap,
                                'visual_capture': visual_capture_path
                            }
                            cell['shapes'].append(component_info)
                            cell['has_visual_content'] = True
                            cell['has_content'] = True
                            
                            # Update display text to show embedded content.
                            # Preserve original cell text whenever it exists — the image
                            # is carried separately in cell['shapes']/['overlapping_shapes']
                            # and rendered as a "Contains: ..." attachment. Only fall back
                            # to a filename/synthetic label when the cell is text-empty.
                            original_text = cell.get('text', '').strip()

                            if component.get('shape_type') == 'Table':
                                cell['priority_component'] = 'Table'
                                if not original_text:
                                    table_id = component.get('box_id', 'unknown')
                                    table_num = table_id[1:] if len(table_id) > 1 else '00'
                                    try:
                                        table_visual = component.get('visual_capture', f"slide_{slide_number:02d}_table_{int(table_num):02d}_visual.png")
                                    except (ValueError, TypeError):
                                        table_visual = component.get('visual_capture', f"slide_{slide_number:02d}_table_{table_num}_visual.png")
                                    cell['display_text'] = f"`{table_visual}`"
                                else:
                                    cell['display_text'] = original_text
                            elif component.get('shape_type') == 'Picture' and cell.get('priority_component') != 'Table':
                                cell['priority_component'] = 'Picture'
                                if not original_text:
                                    picture_id = component.get('box_id', 'unknown')
                                    picture_num = picture_id[1:] if len(picture_id) > 1 else '00'
                                    picture_visual = component.get('visual_capture', f"slide_{slide_number:02d}_picture_{picture_num}_visual.png")
                                    cell['display_text'] = f"`{picture_visual}`"
                                else:
                                    cell['display_text'] = original_text
                            elif component.get('text') and not original_text and not cell.get('priority_component'):
                                cell['display_text'] = component.get('text', '').strip()
                            elif original_text and not cell.get('priority_component'):
                                cell['display_text'] = original_text
                            elif not cell.get('priority_component'):
                                cell['display_text'] = f"{component.get('shape_type')}"
                            
                            # Mark component as consumed only if this is NOT a 1x2 table that will be split
                            table_rows = table_element.get('table_rows', 0)
                            table_cols = table_element.get('table_cols', 0)
                            if not (table_rows == 1 and table_cols == 2):
                                # Only consume for non-1x2 tables (like 3x2 tables)
                                consumed_shape_ids.append(component.get('box_id'))
                                print(f"         🔒 Consumed {component.get('box_id')} (not available for spatial grouping)")
                            else:
                                print(f"         💫 Not consuming {component.get('box_id')} (1x2 table will be split, let spatial grouping handle it)")
                            break
                else:
                    # Check if component has any meaningful overlap with the table overall (even if not with specific cells)
                    table_left = table_element.get('position', {}).get('left', 0)
                    table_top = table_element.get('position', {}).get('top', 0)
                    table_width = table_element.get('position', {}).get('width', 0)
                    table_height = table_element.get('position', {}).get('height', 0)
                    
                    # Calculate overall table overlap
                    table_right = table_left + table_width
                    table_bottom = table_top + table_height
                    
                    intersect_left = max(comp_left, table_left)
                    intersect_top = max(comp_top, table_top)
                    intersect_right = min(comp_left + comp_width, table_right)
                    intersect_bottom = min(comp_top + comp_height, table_bottom)
                    
                    if intersect_left < intersect_right and intersect_top < intersect_bottom:
                        intersect_area = (intersect_right - intersect_left) * (intersect_bottom - intersect_top)
                        comp_area = comp_width * comp_height
                        table_overlap = (intersect_area / comp_area) if comp_area > 0 else 0
                        
                        if table_overlap > 0.50:  # At least 60% overlap with table overall (exclude small overlaps, keep components as standalone groups)
                            print(f"         📎 Adding {component.get('box_id')} as not-embedded component ({table_overlap:.1%} table overlap)")
                            not_embedded_components.append({
                                'box_id': component.get('box_id'),
                                'type': component.get('shape_type'),
                                'text': component.get('text', ''),
                                'table_overlap_percentage': table_overlap,
                                'visual_capture': component.get('visual_capture', '')
                            })
                        else:
                            print(f"         ❌ {component.get('box_id')} has insufficient overlap ({best_overlap:.1%} cell, {table_overlap:.1%} table) - not assigned")
                    else:
                        print(f"         ❌ {component.get('box_id')} has insufficient overlap ({best_overlap:.1%}) - not assigned to any cell")
            
            # ORIGINAL LOGIC (kept for backward compatibility with existing embedded table handling)
            # Require the candidate table to actually sit inside this table's bounds —
            # otherwise sibling tables on the same slide get marked as embedded in each
            # other and both get consumed (then dropped from consolidated_elements).
            parent_left = table_pos.get('left', 0)
            parent_top = table_pos.get('top', 0)
            parent_right = parent_left + table_pos.get('width', 0)
            parent_bottom = parent_top + table_pos.get('height', 0)
            tables_in_overlap = []
            for elem in overlappable_elements:
                if elem.get('shape_type') != 'Table':
                    continue
                ep = elem.get('position', {})
                el = ep.get('left', 0)
                et = ep.get('top', 0)
                er = el + ep.get('width', 0)
                eb = et + ep.get('height', 0)
                eh_w = ep.get('width', 0) * ep.get('height', 0)
                if eh_w <= 0:
                    continue
                ix = max(0, min(er, parent_right) - max(el, parent_left))
                iy = max(0, min(eb, parent_bottom) - max(et, parent_top))
                inside_frac = (ix * iy) / eh_w
                if inside_frac >= 0.5:
                    tables_in_overlap.append(elem)
                else:
                    print(f"         ⏭️  Table {elem.get('box_id')} not embedded in {table_element.get('box_id')} ({inside_frac*100:.1f}% inside)")

            if tables_in_overlap:
                print(f"         🔍 Found {len(tables_in_overlap)} embedded tables - treating as unified components")
                for embedded_table in tables_in_overlap:
                    # Find which cells this embedded table spans
                    table_pos = embedded_table.get('position', {})
                    table_left = table_pos.get('left', 0)
                    table_top = table_pos.get('top', 0)
                    table_width = table_pos.get('width', 0)
                    table_height = table_pos.get('height', 0)
                    
                    # Determine cell coverage for the embedded table
                    start_col = int((table_left - table_left) / cell_width) if cell_width > 0 else 0
                    start_row = int((table_top - table_top) / cell_height) if cell_height > 0 else 2  # Assume bottom rows
                    
                    # Mark cells as containing embedded table
                    for row in range(start_row, min(start_row + 2, table_rows)):  # Span 2 rows typically
                        for col in range(table_cols):
                            for cell in table_element['cell_contents']:
                                if cell['row'] == row and cell['col'] == col:
                                    # Preserve original text if it exists, don't overwrite with generic [TABLE CELL]
                                    if not cell.get('display_text') or cell.get('display_text') == '':
                                        original_text = cell.get('text', '').strip()
                                        cell['display_text'] = original_text if original_text else f'[TABLE CELL {row},{col}]'
                                    cell['has_visual_content'] = True
                                    cell['has_content'] = True
                                    embedded_box_id = embedded_table.get('box_id') or ''
                                    embedded_id_num = embedded_box_id[1:] if embedded_box_id.startswith('S') else embedded_box_id
                                    try:
                                        embedded_id_int = int(embedded_id_num)
                                        embedded_id_part = f"{embedded_id_int:02d}"
                                    except (TypeError, ValueError):
                                        embedded_id_part = embedded_id_num
                                    cell['embedded_table'] = {
                                        'box_id': embedded_box_id,
                                        'visual_capture': f"slide_{slide_number:02d}_table_{embedded_id_part}_visual.png"
                                    }
                                    
                                    # Mark embedded table and its overlapping components as consumed
                                    consumed_shape_ids.append(embedded_table.get('box_id'))
                                    
                    print(f"         ✅ Embedded table {embedded_table.get('box_id')} assigned to cells - will use single table capture")
            
            # Process remaining individual overlapping elements (non-tables)
            individual_elements = [elem for elem in overlappable_elements if elem.get('shape_type') != 'Table']
            for visual_elem in individual_elements:
                visual_pos = visual_elem.get('position', {})
                visual_left = visual_pos.get('left', 0)
                visual_top = visual_pos.get('top', 0)
                visual_width = visual_pos.get('width', 0)
                visual_height = visual_pos.get('height', 0)
                
                # Check if visual element overlaps with table
                if (visual_left < table_left + table_pos.get('width', 0) and
                    visual_left + visual_width > table_left and
                    visual_top < table_top + table_pos.get('height', 0) and
                    visual_top + visual_height > table_top):
                    
                    # Calculate overlap percentage with each table cell to find the best match
                    best_overlap_percentage = 0
                    best_cell_row = 0
                    best_cell_col = 0
                    
                    print(f"         📐 {visual_elem.get('box_id')} calculating overlap percentages...")
                    
                    # Check overlap with each cell using REAL or ESTIMATED boundaries
                    for row in range(table_rows):
                        for col in range(table_cols):
                            if use_real_grid:
                                # Use REAL cell boundaries from XML grid
                                matching_cell = None
                                for grid_cell in real_grid['cells']:
                                    if grid_cell['row'] == row and grid_cell['col'] == col:
                                        matching_cell = grid_cell
                                        break
                                
                                if matching_cell:
                                    cell_left = matching_cell['left']
                                    cell_top = matching_cell['top']
                                    cell_right = matching_cell['right']
                                    cell_bottom = matching_cell['bottom']
                                    print(f"            📐 REAL cell({row},{col}): ({cell_left}, {cell_top}) → ({cell_right}, {cell_bottom})")
                                else:
                                    # Fallback to estimated if specific cell not found
                                    cell_left = table_left + (col * cell_width)
                                    cell_top = table_top + (row * cell_height)
                                    cell_right = cell_left + cell_width
                                    cell_bottom = cell_top + cell_height
                                    print(f"            📐 FALLBACK cell({row},{col}): ({cell_left}, {cell_top}) → ({cell_right}, {cell_bottom})")
                            else:
                                # Use ESTIMATED cell boundaries
                                cell_left = table_left + (col * cell_width)
                                cell_top = table_top + (row * cell_height)
                                cell_right = cell_left + cell_width
                                cell_bottom = cell_top + cell_height
                            
                            # Calculate intersection area
                            intersect_left = max(visual_left, cell_left)
                            intersect_top = max(visual_top, cell_top)
                            intersect_right = min(visual_left + visual_width, cell_right)
                            intersect_bottom = min(visual_top + visual_height, cell_bottom)
                            
                            # Calculate intersection area if there's actual overlap
                            if intersect_left < intersect_right and intersect_top < intersect_bottom:
                                intersect_area = (intersect_right - intersect_left) * (intersect_bottom - intersect_top)
                                element_area = visual_width * visual_height
                                overlap_percentage = (intersect_area / element_area) * 100 if element_area > 0 else 0
                                
                                print(f"            cell({row},{col}): {overlap_percentage:.1f}% overlap")
                                
                                if overlap_percentage > best_overlap_percentage:
                                    best_overlap_percentage = overlap_percentage
                                    best_cell_row = row
                                    best_cell_col = col
                    
                    # Use the SAME threshold as the first pass (cell['shapes']) so a
                    # shape that gets registered into a cell also surfaces in the
                    # rendered Contains: lines. The previous 60% bar dropped images
                    # that had been split across two cells (e.g., the bracket-fixer
                    # photos in slide 6 cell(3,2) at ~51%).
                    second_pass_threshold = 30 if visual_elem.get('shape_type') == 'TextBox' else 50
                    if best_overlap_percentage < second_pass_threshold:
                        print(f"            → No significant overlap found (best: {best_overlap_percentage:.1f}%, need {second_pass_threshold}%), skipping")
                        continue

                    # Redirect away from merged-cell phantoms (same map built earlier).
                    redirect = phantom_to_redirect.get((best_cell_row, best_cell_col))
                    if redirect:
                        print(f"            🔁 Redirecting {visual_elem.get('box_id')} from phantom cell({best_cell_row},{best_cell_col}) → cell{redirect}")
                        best_cell_row, best_cell_col = redirect

                    print(f"            → Best match: cell({best_cell_row},{best_cell_col}) with {best_overlap_percentage:.1f}% overlap")
                    cell_row, cell_col = best_cell_row, best_cell_col
                    
                    # Find the corresponding cell in cell_contents
                    for cell in table_element['cell_contents']:
                        if cell['row'] == cell_row and cell['col'] == cell_col:
                            # Add visual content to this cell
                            if 'overlapping_shapes' not in cell:
                                cell['overlapping_shapes'] = []
                            
                            # Find visual capture path for this shape - use direct filename approach
                            box_id = visual_elem.get('box_id')
                            shape_type = visual_elem.get('shape_type', '').lower()
                            
                            # DEBUG: Log UnifiedGroup detection
                            if box_id and box_id.startswith('G'):
                                print(f"         🔍 DEBUG - UnifiedGroup {box_id}: shape_type='{visual_elem.get('shape_type')}' (lowercase: '{shape_type}')")
                            
                            # Generate expected filename ONLY for capturable shape types
                            # Check if this shape type should be captured (exclude TextBox, AutoShape, etc.)
                            from pptx.enum.shapes import MSO_SHAPE_TYPE
                            shape_type_enum = None
                            if shape_type == 'picture':
                                shape_type_enum = MSO_SHAPE_TYPE.PICTURE
                            elif shape_type == 'table':
                                shape_type_enum = MSO_SHAPE_TYPE.TABLE
                            elif shape_type == 'chart':
                                shape_type_enum = MSO_SHAPE_TYPE.CHART
                            elif shape_type == 'group':
                                shape_type_enum = MSO_SHAPE_TYPE.GROUP
                            elif shape_type == 'embeddedoleobject':
                                shape_type_enum = MSO_SHAPE_TYPE.EMBEDDED_OLE_OBJECT
                            
                            if box_id.startswith('S') and shape_type_enum in self.capture_target_types:
                                try:
                                    shape_index = int(box_id[1:])
                                    if shape_type == 'picture':
                                        expected_filename = f"slide_{slide_number:02d}_picture_{shape_index:02d}_visual.png"
                                    elif shape_type == 'table':
                                        expected_filename = f"slide_{slide_number:02d}_table_{shape_index:02d}_visual.png"
                                    else:
                                        expected_filename = f"slide_{slide_number:02d}_{shape_type}_{shape_index:02d}_visual.png"
                                    
                                    # ASSUME file will be created (don't check existence during overlap detection)
                                    # This fixes timing issue where overlap detection runs before visual capture generation
                                    visual_capture = expected_filename
                                    print(f"         📸 Expected visual capture: {expected_filename}")
                                except (ValueError, IndexError):
                                    visual_capture = None
                            elif box_id.startswith('UG') and shape_type.lower() == 'unifiedgroup':
                                # Handle UnifiedGroup visual captures (UG1, UG2, UG3, etc.)
                                try:
                                    group_index = int(box_id[2:])  # UG1 -> 1, UG2 -> 2
                                    expected_filename = f"slide_{slide_number:02d}_ug{group_index}_visual.png"
                                    # ASSUME file will be created (timing fix)
                                    visual_capture = expected_filename
                                    print(f"         📸 Expected UnifiedGroup visual capture: {expected_filename}")
                                except (ValueError, IndexError):
                                    visual_capture = None
                            elif box_id.startswith('G') and box_id[1:].isdigit() and shape_type.lower() == 'unifiedgroup':
                                # Handle old spatial group visual captures (G1, G2, G3, etc.)
                                try:
                                    group_index = int(box_id[1:])  # G1 -> 1, G2 -> 2
                                    expected_filename = f"slide_{slide_number:02d}_g{group_index}_visual.png"
                                    # ASSUME file will be created (timing fix)
                                    visual_capture = expected_filename
                                    print(f"         📸 Expected spatial group visual capture: {expected_filename}")
                                except (ValueError, IndexError):
                                    visual_capture = None
                            elif box_id.startswith('UG') and not visual_capture:
                                # UnifiedGroup without visual_capture - generate it (fallback)
                                try:
                                    group_index = int(box_id[2:])  # UG1 -> 1, UG2 -> 2
                                    expected_filename = f"slide_{slide_number:02d}_ug{group_index}_visual.png"
                                    visual_capture = expected_filename
                                    print(f"         📸 Generated fallback visual capture for {box_id}: {expected_filename}")
                                except (ValueError, IndexError):
                                    visual_capture = None
                            elif box_id.startswith('G') and box_id[1:].isdigit() and not visual_capture:
                                # Old spatial group without visual_capture - generate it (fallback)
                                try:
                                    group_index = int(box_id[1:])  # G1 -> 1, G2 -> 2
                                    expected_filename = f"slide_{slide_number:02d}_g{group_index}_visual.png"
                                    visual_capture = expected_filename
                                    print(f"         📸 Generated fallback visual capture for {box_id}: {expected_filename}")
                                except (ValueError, IndexError):
                                    visual_capture = None
                            else:
                                # TextBox, AutoShape, etc. - no visual capture
                                visual_capture = None
                                if shape_type in ['textbox', 'autoshape']:
                                    print(f"         📝 Skipping visual capture for {shape_type} (text-only): {box_id}")
                                else:
                                    print(f"         🚫 Skipping visual capture for unsupported type: {shape_type}")
                            
                            # Get overlapping element's text content for TextBox and AutoShape
                            overlapping_text = visual_elem.get('text', '').strip() if visual_elem.get('shape_type') in ['TextBox', 'AutoShape'] else ''
                            
                            cell['overlapping_shapes'].append({
                                'box_id': visual_elem.get('box_id'),
                                'shape_type': visual_elem.get('shape_type'),
                                'overlap_type': 'spatial_overlay',
                                'visual_capture': visual_capture,
                                'overlap_percentage': best_overlap_percentage,
                                'text_content': overlapping_text
                            })
                            cell['has_visual_content'] = True
                            cell['has_content'] = True
                            
                            # Update display text based on element type
                            if not cell.get('text', '').strip():
                                shape_type = visual_elem.get('shape_type', 'Visual')
                                if shape_type == 'Chart':
                                    cell['display_text'] = f"[CHART CELL {cell_row},{cell_col}]"
                                elif shape_type in ['Picture', 'Image']:
                                    cell['display_text'] = f"[IMAGE CELL {cell_row},{cell_col}]"
                                elif shape_type == 'TextBox':
                                    # Show text content for text boxes
                                    preview_text = overlapping_text[:30] + "..." if len(overlapping_text) > 30 else overlapping_text
                                    cell['display_text'] = f"[TEXT CELL {cell_row},{cell_col}]: {preview_text}" if preview_text else f"[TEXT CELL {cell_row},{cell_col}]"
                                elif shape_type == 'AutoShape':
                                    # Show text content for AutoShape headers
                                    if overlapping_text:
                                        preview_text = overlapping_text[:50] + "..." if len(overlapping_text) > 50 else overlapping_text
                                        cell['display_text'] = f"[HEADER {cell_row},{cell_col}]: {preview_text}"
                                    else:
                                        cell['display_text'] = f"[SHAPE CELL {cell_row},{cell_col}]"
                                elif shape_type == 'UnifiedGroup':
                                    # Show UnifiedGroup with its consolidated content info
                                    group_text = visual_elem.get('text', '').strip()
                                    if group_text:
                                        preview_text = group_text[:30] + "..." if len(group_text) > 30 else group_text
                                        cell['display_text'] = f"[GROUP {cell_row},{cell_col}]: {preview_text}"
                                    else:
                                        cell['display_text'] = f"[GROUP CELL {cell_row},{cell_col}]"
                                else:
                                    cell['display_text'] = f"[{shape_type.upper()} CELL {cell_row},{cell_col}]"
                            
                            # Mark this shape as consumed
                            consumed_shape_ids.append(visual_elem.get('box_id'))
                            print(f"      🔗 Associated {visual_elem.get('box_id')} ({visual_elem.get('shape_type')}) with cell({cell_row},{cell_col}) in table {table_element.get('box_id')}")
                            print(f"         ✅ Updated cell({cell_row},{cell_col}): has_content={cell['has_content']}, display_text='{cell['display_text']}'")
                            break
                            
            # Per-cell dedupe: when a Picture and its parent UnifiedGroup both
            # land in the same cell, the group capture is the meaningful one
            # (it includes adjoining arrows / annotations). The individual
            # picture capture was skipped for these members, so a leftover
            # citation would point at a missing file. Drop the Picture entry.
            ug_members = getattr(self, 'unified_group_members', {}) or {}
            if ug_members:
                for cell in table_element.get('cell_contents', []):
                    ug_in_cell = {s.get('box_id') for s in cell.get('overlapping_shapes', [])
                                  if s.get('shape_type') == 'UnifiedGroup'}
                    ug_in_cell |= {s.get('box_id') for s in cell.get('shapes', [])
                                   if s.get('type') == 'UnifiedGroup'}
                    if not ug_in_cell:
                        continue

                    def _is_redundant_member(box_id, shape_type):
                        return (shape_type in ('Picture', 'picture')
                                and ug_members.get(box_id) in ug_in_cell)

                    before_overlap = len(cell.get('overlapping_shapes', []))
                    cell['overlapping_shapes'] = [
                        s for s in cell.get('overlapping_shapes', [])
                        if not _is_redundant_member(s.get('box_id'), s.get('shape_type'))
                    ]
                    cell['shapes'] = [
                        s for s in cell.get('shapes', [])
                        if not _is_redundant_member(s.get('box_id'), s.get('type'))
                    ]
                    dropped = before_overlap - len(cell['overlapping_shapes'])
                    if dropped:
                        print(f"      🧹 Cell({cell.get('row')},{cell.get('col')}): dropped {dropped} redundant Picture citation(s) covered by UnifiedGroup {ug_in_cell}")

        except Exception as e:
            print(f"⚠️  Error associating overlapping shapes with table cells: {e}")

        # Store not-embedded components in the table element for display
        if not_embedded_components:
            table_element['not_embedded_components'] = not_embedded_components
            print(f"      📎 Added {len(not_embedded_components)} not-embedded components to table {table_element.get('box_id')}: {[c.get('box_id') for c in not_embedded_components]}")
        
        # Store consumed shapes globally to prevent individual captures
        if not hasattr(self, 'consumed_by_tables'):
            self.consumed_by_tables = set()
        self.consumed_by_tables.update(consumed_shape_ids)
        
        if consumed_shape_ids:
            print(f"      📋 Tracking {len(consumed_shape_ids)} consumed shapes: {consumed_shape_ids}")
            print(f"      📋 Total consumed shapes now: {list(self.consumed_by_tables)}")
        
        return consumed_shape_ids



    def extract_comprehensive_table_content(self, table_shape):
        """Extract ALL content from a PowerPoint table as a single entity"""
        try:
            table_info = {
                'text': '',
                'has_text': False,
                'table_dimensions': 'unknown',
                'table_rows': 0,
                'table_cols': 0,
                'cell_contents': [],
                'table_type': 'PowerPointTable'
            }
            
            if hasattr(table_shape, 'table'):
                table = table_shape.table
                rows = len(table.rows) if hasattr(table, 'rows') else 0
                cols = len(table.columns) if hasattr(table, 'columns') else 0
                table_info['table_rows'] = rows
                table_info['table_cols'] = cols
                table_info['table_dimensions'] = f"{rows}x{cols}"
                # Extract text from all cells
                all_cell_texts = []
                cell_contents = []
                for row_idx, row in enumerate(table.rows):
                    row_texts = []
                    for col_idx, cell in enumerate(row.cells):
                        cell_text = ''
                        cell_shapes = []
                        
                        # Extract text from cell
                        if hasattr(cell, 'text'):
                            cell_text = cell.text.strip()
                        elif hasattr(cell, 'text_frame') and cell.text_frame:
                            cell_text_parts = []
                            for paragraph in cell.text_frame.paragraphs:
                                for run in paragraph.runs:
                                    cell_text_parts.append(run.text)
                            cell_text = ''.join(cell_text_parts).strip()
                        
                        # Extract shapes (images, charts, etc.) from cell
                        if hasattr(cell, 'text_frame') and cell.text_frame:
                            try:
                                # Look for images
                                for shape in cell.text_frame._element.xpath('.//a:blip'):
                                    cell_shapes.append({'type': 'image', 'element': 'embedded_image'})
                                
                                # Look for charts  
                                for chart in cell.text_frame._element.xpath('.//c:chart', namespaces={'c': 'http://schemas.openxmlformats.org/drawingml/2006/chart'}):
                                    cell_shapes.append({'type': 'chart', 'element': 'embedded_chart'})
                                
                                # Look for other drawing elements (additional visual content)
                                for drawing in cell.text_frame._element.xpath('.//a:graphic'):
                                    cell_shapes.append({'type': 'graphic', 'element': 'embedded_graphic'})
                            except:
                                pass
                        
                        # Check if cell has visual content even without text
                        has_visual_content = len(cell_shapes) > 0
                        
                        # Create cell display text
                        if cell_text:
                            display_text = cell_text
                        elif has_visual_content:
                            # Determine the type of visual content
                            visual_types = [shape['type'] for shape in cell_shapes]
                            if 'chart' in visual_types:
                                display_text = f"[CHART CELL {row_idx},{col_idx}]"
                            elif 'image' in visual_types:
                                display_text = f"[IMAGE CELL {row_idx},{col_idx}]"
                            elif 'graphic' in visual_types:
                                display_text = f"[GRAPHIC CELL {row_idx},{col_idx}]"
                            else:
                                display_text = f"[VISUAL CELL {row_idx},{col_idx}]"
                        else:
                            display_text = ""
                        
                        row_texts.append(display_text)
                        cell_data = {
                            'row': row_idx,
                            'col': col_idx,
                            'text': cell_text,
                            'display_text': display_text,
                            'shapes': cell_shapes,
                            'has_visual_content': has_visual_content,
                            'has_content': bool(cell_text) or has_visual_content
                        }
                        cell_contents.append(cell_data)
                        print(f"      📋 Created cell({row_idx},{col_idx}): text='{cell_text[:20]}...', has_content={cell_data['has_content']}")
                    
                    # Join row cells with separator
                    if any(text for text in row_texts):
                        all_cell_texts.append(' | '.join(row_texts))
                # Combine all table text
                combined_text = ' || '.join(all_cell_texts)
                table_info['text'] = combined_text
                table_info['has_text'] = bool(combined_text)
                table_info['cell_contents'] = cell_contents
            return table_info
            
        except Exception as e:
            print(f"⚠️  Error extracting comprehensive table content: {e}")
            return {
                'text': '',
                'has_text': False,
                'table_dimensions': 'error',
                'table_type': 'PowerPointTable'
            }



    def extract_table_grid_boundaries(self, table_shape):
        """Extract real table grid boundaries from PowerPoint XML structure"""
        try:
            if not hasattr(table_shape, 'table'):
                return None
                
            table = table_shape.table
            table_xml = table_shape._element
            
            # Get table position for absolute coordinates
            table_left = table_shape.left
            table_top = table_shape.top  
            table_width = table_shape.width
            table_height = table_shape.height
            
            print(f"      📐 Extracting real grid from table: pos=({table_left}, {table_top}), size=({table_width}, {table_height})")
            
            # Extract grid information from XML
            grid_info = {
                'table_bounds': {
                    'left': table_left, 'top': table_top,
                    'width': table_width, 'height': table_height
                },
                'rows': len(table.rows),
                'cols': len(table.columns), 
                'cells': []
            }
            
            # Get column widths from table XML
            col_widths = []
            tbl_grid = table_xml.xpath(".//a:tblGrid/a:gridCol")
            if tbl_grid:
                total_grid_width = sum([int(col.get('w', 0)) for col in tbl_grid])
                for col in tbl_grid:
                    col_width_emu = int(col.get('w', 0))
                    # Convert EMU to actual width proportion
                    col_width_actual = (col_width_emu / total_grid_width) * table_width if total_grid_width > 0 else table_width / len(tbl_grid)
                    col_widths.append(col_width_actual)
                print(f"         📏 Found {len(col_widths)} column widths from XML grid")
            else:
                # Fallback: equal column widths
                col_widths = [table_width / len(table.columns) for _ in range(len(table.columns))]
                print(f"         📏 Using equal column widths (no grid found)")
            
            # Calculate actual cell boundaries
            current_top = table_top
            for row_idx, row in enumerate(table.rows):
                row_height = row.height if hasattr(row, 'height') and row.height else table_height / len(table.rows)
                
                current_left = table_left
                for col_idx, cell in enumerate(row.cells):
                    if col_idx < len(col_widths):
                        cell_width = col_widths[col_idx]
                    else:
                        cell_width = table_width / len(table.columns)
                    
                    # Get cell span information from XML
                    cell_xml = table_xml.xpath(f".//a:tbl/a:tr[{row_idx + 1}]/a:tc[{col_idx + 1}]")
                    row_span = 1
                    col_span = 1
                    
                    if cell_xml:
                        cell_element = cell_xml[0]
                        row_span = int(cell_element.get("rowSpan", 1))
                        col_span = int(cell_element.get("gridSpan", 1))
                    
                    # Calculate actual cell boundaries considering spans
                    actual_width = sum(col_widths[col_idx:col_idx + col_span]) if col_idx + col_span <= len(col_widths) else cell_width * col_span
                    actual_height = row_height * row_span
                    
                    cell_bounds = {
                        'row': row_idx, 'col': col_idx,
                        'left': current_left, 'top': current_top,
                        'width': actual_width, 'height': actual_height,
                        'right': current_left + actual_width,
                        'bottom': current_top + actual_height,
                        'row_span': row_span, 'col_span': col_span
                    }
                    
                    grid_info['cells'].append(cell_bounds)
                    current_left += cell_width
                
                current_top += row_height
            
            print(f"         ✅ Extracted real grid: {len(grid_info['cells'])} cells with precise boundaries")
            return grid_info
            
        except Exception as e:
            print(f"         ⚠️  Error extracting table grid: {e}")
            return None



    def check_table_has_overlapping_components(self, table_shape, slide):
        """Check if a table has other components overlapping with its boundary"""
        try:
            # Get table boundary
            table_left = table_shape.left
            table_top = table_shape.top
            table_right = table_shape.left + table_shape.width
            table_bottom = table_shape.top + table_shape.height
            
            # Check all other shapes on the slide
            for other_shape in slide.shapes:
                if other_shape == table_shape:
                    continue  # Skip the table itself
                
                # Skip placeholders and other non-content shapes
                if (hasattr(other_shape, 'shape_type') and 
                    other_shape.shape_type == MSO_SHAPE_TYPE.PLACEHOLDER):
                    continue
                
                # Get other shape boundary
                other_left = other_shape.left
                other_top = other_shape.top
                other_right = other_shape.left + other_shape.width
                other_bottom = other_shape.top + other_shape.height
                
                # Check for overlap (any overlap counts)
                horizontal_overlap = (other_left < table_right and other_right > table_left)
                vertical_overlap = (other_top < table_bottom and other_bottom > table_top)
                
                if horizontal_overlap and vertical_overlap:
                    # Calculate overlap percentage relative to the other shape
                    overlap_left = max(table_left, other_left)
                    overlap_right = min(table_right, other_right)
                    overlap_top = max(table_top, other_top)
                    overlap_bottom = min(table_bottom, other_bottom)
                    
                    overlap_area = (overlap_right - overlap_left) * (overlap_bottom - overlap_top)
                    other_area = other_shape.width * other_shape.height
                    
                    if other_area > 0:
                        overlap_percentage = (overlap_area / other_area) * 100
                        if overlap_percentage > 10:  # 10% overlap threshold
                            print(f"         🔍 Found overlapping component: {other_shape.name if hasattr(other_shape, 'name') else 'Unknown'} ({overlap_percentage:.1f}% overlap)")
                            return True
            
            return False
            
        except Exception as e:
            print(f"      ⚠️  Error checking table overlaps: {e}")
            return False



    def extract_comprehensive_group_content(self, group_shape):
        """Extract ALL content from a PowerPoint group as a single entity"""
        try:
            group_info = {
                'text': '',
                'has_text': False,
                'group_items_count': 0,
                'group_items': [],
                'group_type': 'PowerPointGroup'
            }
            
            if hasattr(group_shape, 'shapes'):
                shapes_in_group = group_shape.shapes
                group_info['group_items_count'] = len(shapes_in_group)
                # Extract content from all shapes in the group
                all_group_texts = []
                group_items = []
                for shape_idx, shape in enumerate(shapes_in_group):
                    shape_text = ''
                    shape_type = self.get_shape_type_name(shape.shape_type)
                    
                    # Extract text from this shape
                    if hasattr(shape, 'text'):
                        shape_text = shape.text.strip()
                    elif hasattr(shape, 'text_frame') and shape.text_frame:
                        text_parts = []
                        for paragraph in shape.text_frame.paragraphs:
                            for run in paragraph.runs:
                                text_parts.append(run.text)
                        shape_text = ''.join(text_parts).strip()
                    
                    # Handle nested tables within groups
                    elif shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                        nested_table_content = self.extract_comprehensive_table_content(shape)
                        shape_text = nested_table_content.get('text', '')
                        shape_type = f"NestedTable({nested_table_content.get('table_dimensions', 'unknown')})"
                    
                    if shape_text:
                        all_group_texts.append(shape_text)
                    
                    group_items.append({
                        'index': shape_idx,
                        'type': shape_type,
                        'text': shape_text,
                        'has_text': bool(shape_text)
                    })
                # Combine all group text
                combined_text = ' | '.join(all_group_texts)
                group_info['text'] = combined_text
                group_info['has_text'] = bool(combined_text)
                group_info['group_items'] = group_items
            return group_info
            
        except Exception as e:
            print(f"⚠️  Error extracting comprehensive group content: {e}")
            return {
                'text': '',
                'has_text': False,
                'group_items_count': 0,
                'group_type': 'PowerPointGroup'
            }



    def consolidate_table_components(self, elements):
        """Apply 75% containment logic only to reasonably-sized PowerPoint tables"""
        powerpoint_tables = [elem for elem in elements if elem.get('table_type') == 'PowerPointTable']
        loose_components = [elem for elem in elements if elem.get('table_type') != 'PowerPointTable']
        
        print(f"      📋 Found {len(powerpoint_tables)} PowerPoint tables - checking reasonable sizes")
        
        # Filter out unreasonably large tables (like S8 that covers 90% of slide)
        reasonable_tables = []
        oversized_tables = []
        
        # Get slide dimensions for percentage calculation
        slide_bounds = self.get_slide_bounds()
        slide_width = slide_bounds['right'] - slide_bounds['left']
        slide_height = slide_bounds['bottom'] - slide_bounds['top']
        
        for table in powerpoint_tables:
            table_width = table['position']['width']
            table_height = table['position']['height']
            
            width_pct = (table_width / slide_width) * 100
            height_pct = (table_height / slide_height) * 100
            
            # Consider table "reasonable" if it covers less than 70% of slide width AND height
            if width_pct < 70.0 and height_pct < 70.0:
                reasonable_tables.append(table)
                print(f"      📋 {table['box_id']}: reasonable size ({width_pct:.1f}%W × {height_pct:.1f}%H) - applying containment")
            else:
                oversized_tables.append(table)
                print(f"      📋 {table['box_id']}: oversized ({width_pct:.1f}%W × {height_pct:.1f}%H) - keeping separate")
        
        # Apply containment logic only to reasonable tables
        enhanced_tables = []
        remaining_components = list(loose_components)
        
        for table in reasonable_tables:
            contained_components = []
            table_pos = table['position']
            
            components_to_remove = []
            for component in remaining_components:
                containment_pct = calculate_spatial_containment(component['position'], table_pos)
                if containment_pct >= 75.0:
                    contained_components.append(component)
                    components_to_remove.append(component)
                    print(f"      📋 {component['box_id']} is {containment_pct:.1f}% contained in {table['box_id']} - merging")
            
            for comp in components_to_remove:
                remaining_components.remove(comp)
            
            if contained_components:
                enhanced_table = self.enhance_powerpoint_table_with_components(table, contained_components)
                enhanced_tables.append(enhanced_table)
                print(f"      📋 Enhanced {table['box_id']} with {len(contained_components)} components")
            else:
                enhanced_tables.append(table)
        
        # Combine: enhanced reasonable tables + oversized tables + remaining loose components
        consolidated_elements = enhanced_tables + oversized_tables + remaining_components
        print(f"      📋 Final: {len(enhanced_tables)} enhanced tables + {len(oversized_tables)} oversized tables + {len(remaining_components)} loose components")
        
        return consolidated_elements



    def is_table_like_structure(self, group):
        """Determine if a group of components forms a table-like structure"""
        if len(group) < 2:  # Reduced from 3 to 2
            return False
        
        # Check if components form a rectangular arrangement
        positions = [comp['position'] for comp in group]
        
        # Get bounding box
        min_left = min(pos['left'] for pos in positions)
        max_right = max(pos['left'] + pos['width'] for pos in positions)
        min_top = min(pos['top'] for pos in positions)
        max_bottom = max(pos['top'] + pos['height'] for pos in positions)
        
        total_width = max_right - min_left
        total_height = max_bottom - min_top
        
        # Check if it's roughly rectangular (not too thin or too wide)
        if total_width == 0 or total_height == 0:
            return False
            
        aspect_ratio = total_width / total_height
        
        # More permissive aspect ratios for table-like structures
        return 0.1 < aspect_ratio < 10.0  # Widened from 0.2-5.0 to 0.1-10.0



    def create_consolidated_table(self, group, table_index):
        """Create a single consolidated table from a group of components"""
        positions = [comp['position'] for comp in group]
        
        # Calculate bounding box
        min_left = min(pos['left'] for pos in positions)
        max_right = max(pos['left'] + pos['width'] for pos in positions)
        min_top = min(pos['top'] for pos in positions)
        max_bottom = max(pos['top'] + pos['height'] for pos in positions)
        
        # Combine text from all components
        text_parts = []
        for comp in group:
            if comp.get('has_text') and comp.get('text'):
                text_parts.append(comp['text'].strip())
        
        combined_text = ' | '.join(text_parts) if text_parts else ''
        
        # Create consolidated table component
        consolidated_table = {
            'box_id': f"T{table_index + 1}",  # T1, T2, T3...
            'shape_index': group[0]['shape_index'],  # Use first component's index
            'position': {
                'left': min_left,
                'top': min_top,
                'width': max_right - min_left,
                'height': max_bottom - min_top
            },
            'shape_type': 'ConsolidatedTable',
            'text': combined_text,
            'has_text': bool(combined_text),
            'box_type': 'table',
            'is_target_for_capture': True,
            'consolidated_components': group,  # Store original components
            'component_count': len(group)
        }
        
        return consolidated_table



    def get_box_type(self, shape_type):
        """Get simplified box type"""
        if shape_type == MSO_SHAPE_TYPE.TEXT_BOX:
            return 'text_box'
        elif shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
            return 'auto_shape'
        elif shape_type == MSO_SHAPE_TYPE.TABLE:
            return 'table'
        elif shape_type == MSO_SHAPE_TYPE.PICTURE:
            return 'picture'
        elif shape_type == MSO_SHAPE_TYPE.CHART:
            return 'chart'
        elif shape_type == MSO_SHAPE_TYPE.GROUP:
            return 'group'
        else:
            return 'other'



    def enhance_powerpoint_table_with_components(self, table, contained_components):
        """Enhance a PowerPoint table by merging contained components into its content"""
        try:
            enhanced_table = table.copy()
            
            # Collect all text content
            table_text = table.get('text', '')
            contained_texts = []
            
            for component in contained_components:
                comp_text = component.get('text', '').strip()
                if comp_text:
                    contained_texts.append(comp_text)
            
            # Merge text content
            if contained_texts:
                if table_text:
                    # Combine table text with contained text using appropriate separator
                    combined_text = table_text + ' || ' + ' | '.join(contained_texts)
                else:
                    combined_text = ' | '.join(contained_texts)
                enhanced_table['text'] = combined_text
                enhanced_table['has_text'] = True
            
            # Add metadata about contained components
            enhanced_table['contained_components'] = contained_components
            enhanced_table['contained_count'] = len(contained_components)
            enhanced_table['enhancement_type'] = 'spatial_containment_75pct'
            
            return enhanced_table
            
        except Exception as e:
            print(f"⚠️  Error enhancing PowerPoint table: {e}")
            return table


