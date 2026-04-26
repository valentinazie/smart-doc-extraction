"""
Slide region screenshots and visual capture utilities.
"""

import os
import time
from pathlib import Path
from datetime import datetime

try:
    from PIL import Image
except ImportError:
    pass


class VisualCaptureMixin:
    """Methods for capturing visual regions from slide images."""

    def should_capture_shape(self, shape, slide, shape_idx):
        """Check if a shape should be captured based on filtering rules"""
        from pptx.enum.shapes import MSO_SHAPE_TYPE
        
        # Skip the individual capture when this picture is already a member of
        # a unified group — the group capture is wider (it includes adjoining
        # arrows / callouts) and the cell renderer is responsible for dropping
        # the redundant individual citation.
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            box_id = f"S{shape_idx}"
            if hasattr(self, 'unified_group_members') and box_id in self.unified_group_members:
                print(f"      🚫 Skipping individual image capture: {box_id} is part of unified group {self.unified_group_members[box_id]}")
                return False

        # Non-table shapes: capture
        if shape.shape_type != MSO_SHAPE_TYPE.TABLE:
            return True
        
        # TABLE FILTERING RULES:
        
        # Rule 0: Skip if this is a boundary table (detected during spatial analysis)
        # Note: We need to find this table in the processed data to check its classification
        box_id = f"S{shape_idx}"
        if hasattr(self, 'comprehensive_data') and 'spatial_analysis' in self.comprehensive_data:
            for slide_data in self.comprehensive_data['spatial_analysis']:
                for element in slide_data.get('boxes', []):
                    if element.get('box_id') == box_id and element.get('is_boundary_table', False):
                        print(f"      🔲 Skipping boundary table capture: {box_id} (table ≈ group size, text-only)")
                        return False
        
        # Rule 1: Skip if table is the group itself (table as container/layout)
        total_shapes = len(slide.shapes)
        if total_shapes <= 2:  # Only slide background + this table = table IS the slide/group
            print(f"      🚫 Skipping table capture: table is the main/only content (layout container)")
            return False
        
        # Rule 2: Skip if table content has no "|" characters (plain text, not actual table)
        try:
            table_text = ""
            if hasattr(shape, 'table') and shape.table:
                # Use SAME extraction logic as main extraction (with | formatting)
                table_rows = []
                for row in shape.table.rows:
                    row_cells = []
                    for cell in row.cells:
                        if cell.text:
                            row_cells.append(cell.text.strip())
                        else:
                            row_cells.append("")
                    if row_cells:  # Only add non-empty rows
                        table_rows.append(" | ".join(row_cells))
                if table_rows:
                    table_text = " | ".join(table_rows)  # Join with | separators like main extraction
                    
            elif hasattr(shape, 'text') and shape.text:
                table_text = shape.text
            
            # Check if it contains table structure indicators
            has_table_structure = "|" in table_text or "\t" in table_text
            if not has_table_structure and table_text.strip():
                print(f"      🚫 Skipping table capture: no table structure found (plain text) - '{table_text[:50]}...'")
                return False
            elif not table_text.strip():
                print(f"      🚫 Skipping table capture: empty table")
                return False
        except Exception as e:
            print(f"      ⚠️  Error analyzing table content: {e}")
            # If we can't analyze, err on the side of caution and capture it
            return True
        
        # Table passed all filters - capture it
        print(f"      ✅ Table has structure - will capture and caption")
        return True



    def capture_target_shapes_both_types(self, slide, slide_number, presentation):
        """Capture both visual (PDF) and shape (matplotlib) captures"""
        print(f"   📸 Capturing target shapes from slide {slide_number}...")
        
        visual_captured = 0
        capture_summary = {}
        
        # Only attempt visual capture if we have slide images from PDF
        if not self.slide_images:
            print(f"      ⚠️  No PDF images available - skipping visual captures")
            print(f"      💡 To enable visual captures, export PPTX as PDF first")
            return
        
        for shape_idx, shape in enumerate(slide.shapes):
            if shape.shape_type in self.capture_target_types:
                # Apply table filtering logic before capture (pass shape_idx for unified group tracking)
                if not self.should_capture_shape(shape, slide, shape_idx):
                    continue
                    
                shape_type_name = self.get_shape_type_name(shape.shape_type)
                # Track what types we're capturing
                if shape_type_name not in capture_summary:
                    capture_summary[shape_type_name] = 0
                # Visual capture from PDF (if available)
                if slide_number in self.slide_images:
                    visual_capture = self.capture_visual_region(shape, slide_number, shape_idx, presentation)
                    if visual_capture:
                        visual_captured += 1
                        capture_summary[shape_type_name] += 1
        
        if visual_captured > 0:
            print(f"      📸 Visual captures: {visual_captured}")
            for shape_type, count in capture_summary.items():
                if count > 0:
                    print(f"         • {shape_type}: {count}")
        else:
            print(f"      📸 No target shapes found for visual capture")



    @staticmethod
    def _rotated_bbox_emu(shape):
        """Return the on-slide axis-aligned bbox (left, top, width, height) of
        a shape, accounting for shape.rotation. python-pptx exposes the
        unrotated bbox; PowerPoint rotates around the rectangle's center."""
        import math

        left = shape.left
        top = shape.top
        width = shape.width
        height = shape.height
        rotation = getattr(shape, 'rotation', 0.0) or 0.0
        # Normalize to [0, 360)
        rotation = rotation % 360
        if rotation == 0:
            return left, top, width, height

        cx = left + width / 2.0
        cy = top + height / 2.0
        theta = math.radians(rotation)
        cos_t = abs(math.cos(theta))
        sin_t = abs(math.sin(theta))
        new_width = width * cos_t + height * sin_t
        new_height = width * sin_t + height * cos_t
        new_left = cx - new_width / 2.0
        new_top = cy - new_height / 2.0
        return int(new_left), int(new_top), int(new_width), int(new_height)


    def capture_visual_region(self, shape, slide_number, shape_index, presentation):
        """Capture the actual visual region from PDF using bounding box coordinates"""
        if slide_number not in self.slide_images:
            return None
        
        try:
            # Get slide image
            slide_image = self.slide_images[slide_number]
            slide_width_px = slide_image.width
            slide_height_px = slide_image.height
            
            # Get shape position in EMU. python-pptx returns the *unrotated*
            # bbox; PowerPoint renders the shape rotated about its center, so
            # for any non-zero rotation we expand the crop to the
            # axis-aligned bounding box of the rotated rectangle. Without this,
            # a 90°/270°-rotated picture (e.g. slide 5 S2 of the Dryer Motor
            # deck) is captured against a tall-narrow box while the PDF shows
            # a wide-short rectangle.
            shape_left_emu, shape_top_emu, shape_width_emu, shape_height_emu = \
                self._rotated_bbox_emu(shape)

            # Get slide dimensions in EMU
            slide_width_emu = presentation.slide_width
            slide_height_emu = presentation.slide_height
            
            # Convert EMU coordinates to pixel coordinates
            left_px = int((shape_left_emu / slide_width_emu) * slide_width_px)
            top_px = int((shape_top_emu / slide_height_emu) * slide_height_px)
            width_px = int((shape_width_emu / slide_width_emu) * slide_width_px)
            height_px = int((shape_height_emu / slide_height_emu) * slide_height_px)
            
            # Add padding
            padding = 10
            left_px = max(0, left_px - padding)
            top_px = max(0, top_px - padding)
            right_px = min(slide_width_px, left_px + width_px + 2*padding)
            bottom_px = min(slide_height_px, top_px + height_px + 2*padding)
            
            # Crop the region
            cropped_region = slide_image.crop((left_px, top_px, right_px, bottom_px))
            
            # Create filename
            shape_type_name = self.get_shape_type_name(shape.shape_type).lower()
            capture_filename = f"slide_{slide_number:02d}_{shape_type_name}_{shape_index:02d}_visual.png"
            capture_filepath = self.visual_captures_dir / capture_filename
            
            # Save the captured region
            cropped_region.save(capture_filepath, 'PNG')
            
            # Create enhanced version with border
            self.create_enhanced_visual_capture(cropped_region, capture_filepath, slide_number, shape_index, shape_type_name)
            
            # Store capture info
            if 'visual_captures' not in self.comprehensive_data:
                self.comprehensive_data['visual_captures'] = []
            
            capture_info = {
                'slide_number': slide_number,
                'shape_index': shape_index,
                'shape_type': shape_type_name,
                'filename': capture_filename,
                'filepath': str(capture_filepath),
                'capture_type': 'visual_pdf',
                'capture_timestamp': datetime.now().isoformat()
            }
            
            self.comprehensive_data['visual_captures'].append(capture_info)
            return capture_info
            
        except Exception as e:
            print(f"⚠️  Error capturing visual region: {e}")
            return None



    def create_enhanced_visual_capture(self, cropped_image, base_filepath, slide_number, shape_index, shape_type):
        """Create enhanced version with border and information"""
        try:
            enhanced_filename = f"slide_{slide_number:02d}_{shape_type}_{shape_index:02d}_enhanced.png"
            enhanced_filepath = base_filepath.parent / enhanced_filename
            
            # Add border and text
            border_size = 40
            enhanced_width = cropped_image.width + 2 * border_size
            enhanced_height = cropped_image.height + 2 * border_size
            
            # Create new image with border
            enhanced_image = Image.new('RGB', (enhanced_width, enhanced_height), 'white')
            enhanced_image.paste(cropped_image, (border_size, border_size))
            
            # Add border rectangle
            draw = ImageDraw.Draw(enhanced_image)
            draw.rectangle([border_size-2, border_size-2, enhanced_width-border_size+1, enhanced_height-border_size+1], 
                          outline='red', width=3)
            
            # Add text information
            try:
                font = ImageFont.truetype("Arial.ttf", 14)
            except:
                font = ImageFont.load_default()
            
            info_text = f"Slide {slide_number} - {shape_type.title()} {shape_index} (Comprehensive)"
            draw.text((10, 10), info_text, fill='black', font=font)
            
            enhanced_image.save(enhanced_filepath, 'PNG')
            
        except Exception as e:
            print(f"⚠️  Could not create enhanced capture: {e}")



    def find_visual_capture_file(self, box_id, slide_number):
        """Find the visual capture file for a given box ID and slide number"""
        
        # Handle Image+AutoShape consolidated entities (IA prefix - simple script style)
        if box_id.startswith('IA') and not box_id.startswith('IA_'):
            # Simple script style: IA3, IA11, etc.
            # Look for the new consolidated entity capture
            consolidated_filename = f"slide_{slide_number:02d}_{box_id}_visual.png"
            consolidated_path = self.visual_captures_dir / consolidated_filename
            if consolidated_path.exists():
                return consolidated_filename
            return f"→ {box_id} (consolidated image+autoshape combo)"
        
        # Handle Image+AutoShape consolidated entities (IA_ prefix - old complex style)
        if box_id.startswith('IA_'):
            # Try individual component captures first (from new capture method)
            individual_patterns = [
                f"slide_{slide_number:02d}_individual_{box_id}_primary_*.png",
                f"slide_{slide_number:02d}_individual_{box_id}_*.png"
            ]
            
            for pattern in individual_patterns:
                import glob
                matches = glob.glob(str(self.visual_captures_dir / pattern))
                if matches:
                    return Path(matches[0]).name
            return f"→ {box_id} (consolidated image+autoshape combo)"
        
        # Handle consolidated images (CI prefix)
        if box_id.startswith('CI'):
            # Consolidated image naming: slide_{slide_number:02d}_consolidated_image_{box_id}.png
            consolidated_filename = f"slide_{slide_number:02d}_consolidated_image_{box_id}.png"
            consolidated_path = self.visual_captures_dir / consolidated_filename
            if consolidated_path.exists():
                return consolidated_filename
            return f"→ {box_id} (consolidated image)"
        
        # Handle regular shapes (original logic)
        visual_captures = self.comprehensive_data.get('visual_captures', [])
        # Debug: print(f"      📸 DEBUG - find_visual_capture_file({box_id}, {slide_number}): {len(visual_captures)} captures available")
        
        # First try direct box_id match (for consolidated entities)
        for capture in visual_captures:
            if (capture['slide_number'] == slide_number and 
                capture.get('box_id') == box_id):
                return capture['filename']
        
        # Extract shape index from box_id (e.g., "S18" -> 18)
        try:
            if box_id.startswith('S'):
                shape_index = int(box_id[1:])
            else:
                shape_index = int(box_id)
        except (ValueError, IndexError):
            return None
        
        # Find matching visual capture by shape_index and ensure it's the right type
        # For tables (S-prefix), look for table-type captures
        if box_id.startswith('S'):
            for capture in visual_captures:
                if (capture['slide_number'] == slide_number and 
                    capture.get('shape_index') == shape_index and
                    capture.get('shape_type', '').lower() == 'table'):
                    return capture['filename']
        
        # For non-table components, use original logic
        for capture in visual_captures:
            if (capture['slide_number'] == slide_number and 
                capture.get('shape_index') == shape_index):
                return capture['filename']
        
        return None



    def capture_hybrid_entity_images(self, hybrid_entity, slide_number):
        """Capture images referenced by hybrid entity while preserving text visibility"""
        if slide_number not in self.slide_images:
            print(f"      ⚠️  No slide image available for slide {slide_number} - skipping hybrid capture")
            return None
        
        hybrid_images = hybrid_entity.get('hybrid_images', [])
        if not hybrid_images:
            return None
        
        print(f"      📸 Capturing hybrid entity {hybrid_entity['box_id']} (text preserved, {len(hybrid_images)} images captured)")
        
        capture_count = 0
        
        for image_id in hybrid_images:
            # Create a capture filename for this hybrid image
            capture_filename = f"slide_{slide_number:02d}_hybrid_{hybrid_entity['box_id']}_{image_id}_visual.png"
            print(f"         📸 Hybrid image capture: {capture_filename}")
            capture_count += 1
        
        return capture_count > 0



    def capture_consolidated_image_visual(self, consolidated_image, slide_number):
        """Capture visual region for a ConsolidatedImage entity using its combined boundary"""
        if slide_number not in self.slide_images:
            return None
        
        # Get the combined boundary position
        position = consolidated_image['position']
        box_id = consolidated_image['box_id']
        
        # Use the same visual capture logic but with the consolidated boundary
        try:
            slide_image = self.slide_images[slide_number]
            
            # Convert PowerPoint coordinates to image pixels
            # PowerPoint uses EMUs (English Metric Units), we need to convert to pixels
            # Slide dimensions in PowerPoint: typically 10" x 7.5" = 9144000 x 6858000 EMUs
            # PDF image dimensions vary based on DPI
            img_width, img_height = slide_image.size
            
            # PowerPoint slide dimensions (standard 16:9)
            # Get actual slide dimensions from metadata
            ppt_width = self.comprehensive_data.get("metadata", {}).get("slide_width", 9144000)
            ppt_height = self.comprehensive_data.get("metadata", {}).get("slide_height", 6858000)
            
            # Convert coordinates
            x_scale = img_width / ppt_width
            y_scale = img_height / ppt_height
            
            left = int(position['left'] * x_scale)
            top = int(position['top'] * y_scale)
            width = int(position['width'] * x_scale)
            height = int(position['height'] * y_scale)
            
            # Ensure bounds are within image
            left = max(0, left)
            top = max(0, top)
            right = min(img_width, left + width)
            bottom = min(img_height, top + height)
            
            if right <= left or bottom <= top:
                print(f"      ⚠️  Invalid crop bounds for {box_id}")
                return None
            
            # Crop the consolidated region
            cropped_region = slide_image.crop((left, top, right, bottom))
            
            # Add border to make it more visible
            bordered_region = ImageOps.expand(cropped_region, border=3, fill='red')
            
            # Save the consolidated image capture
            capture_filename = f"slide_{slide_number:02d}_consolidated_image_{box_id}.png"
            capture_path = self.visual_captures_dir / capture_filename
            bordered_region.save(capture_path, 'PNG', dpi=(600, 600))
            
            print(f"      🖼️  Consolidated image captured: {capture_filename}")
            return capture_filename
            
        except Exception as e:
            print(f"      ❌ Failed to capture consolidated image {box_id}: {e}")
            return None



    def capture_individual_components_of_consolidated_entity(self, consolidated_entity, slide_number):
        """Capture individual components of a consolidated entity to avoid text inclusion"""
        if slide_number not in self.slide_images:
            return None
        
        box_id = consolidated_entity['box_id']
        box_type = consolidated_entity['box_type']
        
        if box_type == 'consolidated_image':
            # For consolidated images, capture each constituent image separately
            constituent_images = consolidated_entity.get('constituent_images', [])
            print(f"      🖼️  Capturing {len(constituent_images)} individual images in consolidated entity {box_id}")
            
            for img_info in constituent_images:
                img_position = img_info['original_position']
                img_id = img_info['box_id']
                self.capture_region_from_position(img_position, slide_number, f"{box_id}_{img_id}")
        elif box_type == 'image_autoshape_combo':
            # For image+autoshape combos, only capture the primary image, not the autoshapes
            primary_image = consolidated_entity.get('primary_image', {})
            if primary_image:
                img_position = primary_image['original_position']
                img_id = primary_image['box_id']
                print(f"      🖼️  Capturing primary image {img_id} from combo {box_id} (excluding AutoShapes to avoid text)")
                self.capture_region_from_position(img_position, slide_number, f"{box_id}_primary_{img_id}")
            
            # Optionally capture AutoShapes only if they're actually visual (not text containers)
            autoshapes = consolidated_entity.get('overlapping_autoshapes', [])
            for autoshape_info in autoshapes:
                autoshape_id = autoshape_info['box_id']
                autoshape_position = autoshape_info['original_position']
                # Only capture if it's likely a visual element (has significant size)
                if autoshape_position['width'] > 50000 and autoshape_position['height'] > 50000:  # ~4pts minimum
                    print(f"      🔧 Capturing visual AutoShape {autoshape_id} from combo {box_id}")
                    self.capture_region_from_position(autoshape_position, slide_number, f"{box_id}_autoshape_{autoshape_id}")



    def capture_region_from_position(self, position, slide_number, capture_id):
        """Capture a visual region from a position dictionary"""
        try:
            slide_image = self.slide_images[slide_number]
            
            # Convert PowerPoint coordinates to image pixels
            img_width, img_height = slide_image.size
            # Get actual slide dimensions from metadata
            ppt_width = self.comprehensive_data.get("metadata", {}).get("slide_width", 9144000)
            ppt_height = self.comprehensive_data.get("metadata", {}).get("slide_height", 6858000)
            
            # Convert coordinates
            x_scale = img_width / ppt_width
            y_scale = img_height / ppt_height
            
            left = int(position['left'] * x_scale)
            top = int(position['top'] * y_scale)
            width = int(position['width'] * x_scale)
            height = int(position['height'] * y_scale)
            
            # Ensure bounds are within image
            left = max(0, left)
            top = max(0, top)
            right = min(img_width, left + width)
            bottom = min(img_height, top + height)
            
            if right <= left or bottom <= top:
                print(f"      ⚠️  Invalid crop bounds for {capture_id}")
                return None
            
            # Crop the region
            cropped_region = slide_image.crop((left, top, right, bottom))
            
            # Add subtle border to distinguish from combined captures
            bordered_region = ImageOps.expand(cropped_region, border=2, fill='blue')
            
            # Save individual component capture
            capture_filename = f"slide_{slide_number:02d}_individual_{capture_id}.png"
            capture_path = self.visual_captures_dir / capture_filename
            bordered_region.save(capture_path, 'PNG', dpi=(600, 600))
            
            print(f"      📸 Individual component captured: {capture_filename}")
            return capture_filename
            
        except Exception as e:
            print(f"      ❌ Failed to capture individual component {capture_id}: {e}")
            return None



    def capture_consolidated_entity_visual(self, consolidated_entity, slide_number):
        """Capture visual region for consolidated entities (IA~ groups)"""
        if slide_number not in self.slide_images:
            print(f"      ⚠️  No slide image available for slide {slide_number} - skipping IA~ capture")
            return None
        
        box_id = consolidated_entity['box_id']
        box_type = consolidated_entity.get('box_type', 'unknown')
        
        print(f"      📸 Capturing consolidated entity {box_id} ({box_type})")
        
        try:
            slide_image = self.slide_images[slide_number]
            position = consolidated_entity['position']
            
            # Convert PowerPoint coordinates to image pixels
            img_width, img_height = slide_image.size
            # Get actual slide dimensions from metadata
            ppt_width = self.comprehensive_data.get("metadata", {}).get("slide_width", 9144000)
            ppt_height = self.comprehensive_data.get("metadata", {}).get("slide_height", 6858000)
            
            # Convert coordinates
            x_scale = img_width / ppt_width
            y_scale = img_height / ppt_height
            
            left = int(position['left'] * x_scale)
            top = int(position['top'] * y_scale)
            width = int(position['width'] * x_scale)
            height = int(position['height'] * y_scale)
            
            # Add padding
            padding = 10
            left = max(0, left - padding)
            top = max(0, top - padding)
            right = min(img_width, left + width + 2*padding)
            bottom = min(img_height, top + height + 2*padding)
            
            # Ensure valid bounds
            if right <= left or bottom <= top:
                print(f"      ⚠️  Invalid crop bounds for {box_id}")
                return None
            
            # Crop the consolidated region
            cropped_region = slide_image.crop((left, top, right, bottom))
            
            # Add border to distinguish consolidated captures
            bordered_region = ImageOps.expand(cropped_region, border=4, fill='blue')
            
            # Save the consolidated entity capture
            capture_filename = f"slide_{slide_number:02d}_{box_id}_visual.png"
            capture_path = self.visual_captures_dir / capture_filename
            bordered_region.save(capture_path, 'PNG', dpi=(600, 600))
            
            # Add to comprehensive data visual captures list
            capture_info = {
                'slide_number': slide_number,
                'box_id': box_id,
                'box_type': box_type,
                'shape_index': consolidated_entity.get('shape_index', 0),  # Include shape_index for consistency
                'filename': capture_filename,
                'filepath': str(capture_path),
                'entity_type': 'consolidated',
                'capture_timestamp': datetime.now().isoformat()
            }
            
            if 'visual_captures' not in self.comprehensive_data:
                self.comprehensive_data['visual_captures'] = []
            self.comprehensive_data['visual_captures'].append(capture_info)
            
            print(f"      🖼️  Captured consolidated entity: {capture_filename}")
            return capture_filename
            
        except Exception as e:
            print(f"      ❌ Failed to capture consolidated entity {box_id}: {e}")
            return None


