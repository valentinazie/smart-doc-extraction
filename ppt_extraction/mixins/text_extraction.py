"""
Text extraction via Watsonx API and native python-pptx.
"""

import os
import json
import time
from pathlib import Path
from datetime import datetime

try:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

try:
    from ibm_watsonx_ai.foundation_models.extractions import TextExtractionsV2, TextExtractionsV2ResultFormats
    from ibm_watsonx_ai.helpers import DataConnection, S3Location
    from ibm_watsonx_ai.metanames import TextExtractionsV2ParametersMetaNames
    WATSONX_SDK_AVAILABLE = True
except ImportError:
    WATSONX_SDK_AVAILABLE = False


class TextExtractionMixin:
    """Methods for extracting text content from presentations."""

    def extract_text_structure(self, pptx_path):
        """Extract hierarchical text structure using watsonx API or fallback to native"""
        print(f"\n📝 EXTRACTING TEXT STRUCTURE")
        print("=" * 50)
        
        if self.watsonx_available:
            try:
                print("🚀 Using watsonx Text Extraction V2 API...")
                
                # Step 1: Convert PPTX to PDF for watsonx
                pdf_path = self.convert_pptx_to_pdf_for_watsonx(pptx_path)
                if not pdf_path:
                    print("🔄 PDF conversion failed, falling back to native extraction")
                    return self.extract_text_structure_native(pptx_path)
                
                # Step 2: Extract text with watsonx API
                watsonx_results = self.extract_text_with_watsonx_api(pdf_path)
                if not watsonx_results:
                    print("🔄 watsonx extraction failed, falling back to native extraction")
                    return self.extract_text_structure_native(pptx_path)
                
                # Step 3: Extract spatial structure from PPTX
                print("   📦 Extracting spatial structure from PPTX...")
                spatial_structure = self.extract_spatial_structure_native(pptx_path)
                if not spatial_structure:
                    print("🔄 Spatial extraction failed, falling back to native extraction")
                    return self.extract_text_structure_native(pptx_path)
                
                # Step 4: Extract tables from watsonx results
                print("   📊 Extracting tables from watsonx results...")
                table_data = self.extract_tables_from_watsonx(watsonx_results)
                table_files = []
                if table_data:
                    table_files = self.save_table_extractions(table_data, 1)
                
                # Step 5: Combine watsonx text with spatial positioning
                combined_structure = self.combine_watsonx_text_with_spatial(watsonx_results, spatial_structure)
                
                # Step 6: Add table references to combined structure
                if table_files:
                    combined_structure['table_extractions'] = {
                        'enabled': True,
                        'table_files': table_files,
                        'total_tables': len(table_data)
                    }
                
                # Step 7: Save results
                self.comprehensive_data['text_structure'] = combined_structure
                text_structure_path = self.text_structure_dir / "text_structure.json"
                with open(text_structure_path, 'w', encoding='utf-8') as f:
                    json.dump(combined_structure, f, ensure_ascii=False, indent=2)
                
                # Export slides as individual text files
                self.export_slides_as_text(combined_structure)
                
                print(f"✅ Text structure extracted with watsonx: {len(combined_structure['slides'])} slides")
                return True
                
            except Exception as e:
                print(f"❌ Error with watsonx text extraction: {e}")
                print("🔄 Falling back to native extraction")
                return self.extract_text_structure_native(pptx_path)
        else:
            print("⚠️  watsonx not available - using native python-pptx extraction")
            return self.extract_text_structure_native(pptx_path)

            return self.extract_text_structure_native(pptx_path)



    def extract_text_structure_native(self, pptx_path):
        """Extract hierarchical text structure using native python-pptx (fallback)"""
        if not PPTX_AVAILABLE:
            print("❌ python-pptx not available")
            return False
        
        try:
            prs = Presentation(pptx_path)
            
            text_structure = {
                'file_info': {
                    'name': Path(pptx_path).name,
                    'path': str(Path(pptx_path).absolute()),
                    'total_slides': len(prs.slides),
                    'extraction_method': 'native_python_pptx',
                },
                'slides': []
            }
            
            if not prs.slides:
                print("   ❌ No slides found in PPTX")
                return False
            for slide_num, slide in enumerate(prs.slides, 1):
                print(f"   📄 Processing slide {slide_num}/{len(prs.slides)} text structure...")
                slide_structure = self.extract_slide_text_structure(slide, slide_num)
                text_structure['slides'].append(slide_structure)
            
            self.comprehensive_data['text_structure'] = text_structure
            
            # Save text structure
            text_structure_path = self.text_structure_dir / "text_structure.json"
            with open(text_structure_path, 'w', encoding='utf-8') as f:
                json.dump(text_structure, f, ensure_ascii=False, indent=2)
            
            # Export slides as individual text files
            self.export_slides_as_text(text_structure)
            
            print(f"✅ Text structure extracted: {len(prs.slides)} slides")
            return True
            
        except Exception as e:
            print(f"❌ Error extracting text structure: {e}")
            return False



    def extract_text_with_watsonx_api(self, pdf_path):
        """Extract text from PDF using watsonx Text Extraction V2 API"""
        try:
            pdf_path = Path(pdf_path)
            timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
            
            # Upload PDF to COS
            input_file = f"input/{pdf_path.stem}_{timestamp}.pdf"
            output_directory = f"output/{pdf_path.stem}_{timestamp}/"  # Directory with trailing slash
            
            print(f"   📤 Uploading PDF to COS: {input_file}")
            
            # Upload file using COS client (correct method from test file)
            self.cos_client.upload_file(str(pdf_path), self.cos_bucket_name, input_file)
            print(f"   📤 Uploaded successfully")
            
            input_ref = DataConnection(
                connection_asset_id=self.cos_connection_id,
                location=S3Location(bucket=self.cos_bucket_name, path=input_file)
            )
            input_ref.set_client(self.watsonx_client)  # This should be self.watsonx_client (our APIClient)
            
            # Setup output reference
            output_ref = DataConnection(
                connection_asset_id=self.cos_connection_id,
                location=S3Location(bucket=self.cos_bucket_name, path=output_directory)
            )
            output_ref.set_client(self.watsonx_client)
            
            # Create watsonx extraction job
            extraction = TextExtractionsV2(credentials=self.credentials, space_id=self.space_id)
            
            print(f"   🚀 Starting watsonx text extraction job...")
            response = extraction.run_job(
                document_reference=input_ref,
                results_reference=output_ref,
                result_formats=[
                    TextExtractionsV2ResultFormats.ASSEMBLY_JSON,
                    TextExtractionsV2ResultFormats.MARKDOWN
                ],
                parameters={
                    TextExtractionsV2ParametersMetaNames.MODE: "high_quality",
                    TextExtractionsV2ParametersMetaNames.OCR_MODE: "enabled",
                    TextExtractionsV2ParametersMetaNames.LANGUAGES: ["ko", "en"],
                    TextExtractionsV2ParametersMetaNames.OUTPUT_DPI: 1200,  # Higher DPI like working script
                    TextExtractionsV2ParametersMetaNames.AUTO_ROTATION_CORRECTION: True,  # Handle orientation
                    TextExtractionsV2ParametersMetaNames.CREATE_EMBEDDED_IMAGES: "disabled",  # Image processing
                    TextExtractionsV2ParametersMetaNames.OUTPUT_TOKENS_AND_BBOX: True,
                }
            )
            
            job_id = TextExtractionsV2.get_job_id(response)
            print(f"   ⏳ watsonx job started - ID: {job_id}")
            
            # Wait for completion
            import time
            max_wait = 300  # 5 minutes max
            wait_time = 0
            poll_interval = 10
            
            while wait_time < max_wait:
                time.sleep(poll_interval)
                wait_time += poll_interval
                try:
                    job_details = extraction.get_job_details(job_id)
                    status = job_details['entity']['results']['status']
                    print(f"   📊 Status after {wait_time}s: {status}")
                    
                    if status == 'completed':
                        # Download results - discover assembly.json in output directory
                        print(f"   📥 Downloading watsonx results...")
                        
                        # List files in output directory to find assembly.json and markdown
                        try:
                            objects = self.cos_client.list_objects_v2(
                                Bucket=self.cos_bucket_name,
                                Prefix=output_directory
                            )
                            
                            assembly_file = None
                            markdown_file = None
                            if 'Contents' in objects:
                                for obj in objects['Contents']:
                                    if obj['Key'].endswith('assembly.json'):
                                        assembly_file = obj['Key']
                                    elif obj['Key'].endswith('assembly.md') or obj['Key'].endswith('.md'):
                                        markdown_file = obj['Key']
                            
                            results_json = None
                            saved_files = []
                            
                            # Download and save assembly.json (for processing and as raw output)
                            if assembly_file:
                                import tempfile
                                with tempfile.NamedTemporaryFile(mode='w+b', delete=False) as temp_file:
                                    self.cos_client.download_fileobj(self.cos_bucket_name, assembly_file, temp_file)
                                    temp_file.seek(0)
                                    
                                with open(temp_file.name, 'r', encoding='utf-8') as f:
                                    results_json = json.load(f)
                                
                                # Save raw assembly.json to watsonx_raw_outputs_dir
                                assembly_output_path = self.watsonx_raw_outputs_dir / f"{pdf_path.stem}_comprehensive_assembly.json"
                                with open(assembly_output_path, 'w', encoding='utf-8') as f:
                                    json.dump(results_json, f, ensure_ascii=False, indent=2, default=str)
                                saved_files.append(f"assembly: {assembly_output_path.name}")
                                
                                import os
                                os.unlink(temp_file.name)
                            
                            # Download and save markdown file
                            if markdown_file:
                                markdown_content = ""
                                try:
                                    import tempfile
                                    with tempfile.NamedTemporaryFile(mode='w+b', delete=False) as temp_file:
                                        self.cos_client.download_fileobj(self.cos_bucket_name, markdown_file, temp_file)
                                        temp_file.seek(0)
                                        
                                    with open(temp_file.name, 'r', encoding='utf-8') as f:
                                        markdown_content = f.read()
                                    
                                    # Save raw markdown to watsonx_raw_outputs_dir
                                    markdown_output_path = self.watsonx_raw_outputs_dir / f"{pdf_path.stem}_comprehensive_markdown.md"
                                    with open(markdown_output_path, 'w', encoding='utf-8') as f:
                                        f.write(markdown_content)
                                    saved_files.append(f"markdown: {markdown_output_path.name}")
                                    
                                    import os
                                    os.unlink(temp_file.name)
                                except Exception as e:
                                    print(f"   ⚠️  Warning: Could not download markdown file: {e}")
                            
                            if results_json:
                                print(f"   ✅ watsonx text extraction completed successfully")
                                print(f"   🤖 CONFIRMED: Using watsonx Text Extraction V2 API for high-quality text extraction")
                                if saved_files:
                                    print(f"   💾 Saved raw watsonx outputs: {', '.join(saved_files)}")
                                return results_json
                            else:
                                print(f"   ❌ No assembly.json found in output directory")
                                return None
                                
                        except Exception as download_error:
                            print(f"   ❌ Error downloading results: {download_error}")
                            return None
                    elif status == 'failed':
                        print(f"   ❌ watsonx job failed")
                        return None
                        
                except Exception as e:
                    print(f"   ⚠️  Status check error: {e}")
            
            print(f"   ⏰ watsonx job timeout after {max_wait}s")
            return None
            
        except Exception as e:
            print(f"   ❌ watsonx API error: {e}")
            return None



    def combine_watsonx_text_with_spatial(self, watsonx_results, spatial_structure):
        """Combine watsonx extracted text with spatial positioning from PPTX"""
        try:
            # Extract all text from watsonx assembly results
            watsonx_text_by_page = self.extract_text_by_page_from_watsonx(watsonx_results)
            
            combined_structure = {
                'file_info': spatial_structure['file_info'].copy(),
                'slides': []
            }
            combined_structure['file_info']['extraction_method'] = 'watsonx_with_spatial_mapping'
            
            # Combine text and spatial data slide by slide
            for slide_idx, spatial_slide in enumerate(spatial_structure['slides']):
                slide_num = spatial_slide['slide_number']
                # Get watsonx text items for this page/slide (now an array of individual items)
                watsonx_text_items = watsonx_text_by_page.get(slide_num - 1, [])  # 0-based indexing
                
                combined_slide = {
                    'slide_number': slide_num,
                    'title': f'Slide {slide_num}',
                    'content': [],
                    'notes': '',
                    'layout_name': 'Unknown',
                    'watsonx_text_items': watsonx_text_items,  # Store individual watsonx text items
                    'spatial_shapes': spatial_slide['shapes']
                }
                
                # Create content entries from individual watsonx text items (no distribution needed!)
                for idx, text_item in enumerate(watsonx_text_items):
                    if text_item.strip():
                        content_item = {
                            'content_id': f'slide_{slide_num}_content_{idx + 1}',
                            'text': text_item,  # Use individual text item directly
                            'source': 'watsonx_api',
                            'confidence': True,  # watsonx is considered high confidence
                            'is_title': False,   
                            'reading_order': idx + 1  
                        }
                        combined_slide['content'].append(content_item)
                combined_structure['slides'].append(combined_slide)
            
            return combined_structure
            
        except Exception as e:
            print(f"   ❌ Error combining watsonx and spatial data: {e}")
            return None



    def extract_text_by_page_from_watsonx(self, watsonx_results):
        """Extract text organized by page from watsonx assembly results using correct structure"""
        text_by_page = {}
        
        # DEBUG: Show watsonx results structure to find the right level
        print(f"   🔍 DEBUG - watsonx_results keys: {list(watsonx_results.keys())}")
        if 'all_structures' in watsonx_results:
            print(f"   🔍 DEBUG - all_structures keys: {list(watsonx_results['all_structures'].keys())}")
        
        # Check for other structured elements that might have meaningful text chunks
        for key in ['pages', 'paragraphs', 'text_blocks', 'sections', 'sentences']:
            if key in watsonx_results:
                count = len(watsonx_results[key]) if isinstance(watsonx_results[key], list) else 1
                print(f"   🔍 DEBUG - Found {key}: {count} items")
                if isinstance(watsonx_results[key], list) and len(watsonx_results[key]) > 0:
                    sample = watsonx_results[key][0]
                    if isinstance(sample, dict) and 'text' in sample:
                        print(f"      📝 Sample {key}[0]: {repr(sample['text'][:100])}")
        
        try:
            # FIXED: Access structured elements (paragraphs, text_blocks) instead of concatenated tokens
            structured_text_items = []
            
            # Method 1: Try section_titles from all_structures (HIGHEST PRIORITY - clean formatted text!)
            if 'all_structures' in watsonx_results and 'section_titles' in watsonx_results['all_structures']:
                for title in watsonx_results['all_structures']['section_titles']:
                    if 'text' in title and title['text'].strip():
                        structured_text_items.append(title['text'].strip())
                print(f"   📄 Extracted {len(structured_text_items)} section_titles from watsonx all_structures (CLEAN TEXT!)")
            
            # Method 2: Try list_items from all_structures (good for bullet points)
            elif 'all_structures' in watsonx_results and 'list_items' in watsonx_results['all_structures']:
                for item in watsonx_results['all_structures']['list_items']:
                    if 'text' in item and item['text'].strip():
                        structured_text_items.append(item['text'].strip())
                print(f"   📄 Extracted {len(structured_text_items)} list_items from watsonx all_structures")
            
            # Method 3: Try paragraphs from all_structures (fallback)
            elif 'all_structures' in watsonx_results and 'paragraphs' in watsonx_results['all_structures']:
                for para in watsonx_results['all_structures']['paragraphs']:
                    if 'text' in para and para['text'].strip():
                        structured_text_items.append(para['text'].strip())
                print(f"   📄 Extracted {len(structured_text_items)} paragraphs from watsonx all_structures")
            
            # Method 4: Try sections from all_structures
            elif 'all_structures' in watsonx_results and 'sections' in watsonx_results['all_structures']:
                for section in watsonx_results['all_structures']['sections']:
                    if 'text' in section and section['text'].strip():
                        structured_text_items.append(section['text'].strip())
                print(f"   📄 Extracted {len(structured_text_items)} sections from watsonx all_structures")
            
            # Method 5: Try all_structures.tokens but keep individual items (old fallback)
            elif 'all_structures' in watsonx_results and 'tokens' in watsonx_results['all_structures']:
                tokens = watsonx_results['all_structures']['tokens']
                # IMPORTANT: Keep individual token texts, don't concatenate
                for item in tokens:
                    if 'text' in item and item['text'].strip():
                        structured_text_items.append(item['text'].strip())
                print(f"   📄 Extracted {len(structured_text_items)} individual tokens from watsonx assembly")
            
            if structured_text_items:
                # Store individual text items for page 0 (slide 1)
                text_by_page[0] = structured_text_items  # Array of individual items, not concatenated!
                return text_by_page
            else:
                print(f"   ⚠️  No structured text found in watsonx results, using fallback extraction")
        except Exception as e:
            print(f"   ⚠️  Error in primary extraction: {e}, using fallback")
        
        # Fallback: original recursive method
        def extract_text_recursive(obj):
            if isinstance(obj, dict):
                # Check if this object has page information and text
                if 'page' in obj and 'text' in obj:
                    page_num = obj['page']
                    text = obj['text']
                    if page_num not in text_by_page:
                        text_by_page[page_num] = []
                    text_by_page[page_num].append(text)
                # Recursively search all values
                for key, value in obj.items():
                    if key in ['text', 'raw_text', 'content'] and isinstance(value, str):
                        # Assume page 0 if no page info found
                        if 0 not in text_by_page:
                            text_by_page[0] = []
                        text_by_page[0].append(value)
                    elif isinstance(value, (dict, list)):
                        extract_text_recursive(value)
            elif isinstance(obj, list):
                for item in obj:
                    extract_text_recursive(item)
        
        extract_text_recursive(watsonx_results)
        
        # Keep lists as individual items (don't concatenate)
        # text_by_page now contains arrays of individual text items
        
        return text_by_page



    def distribute_watsonx_text_to_shapes(self, watsonx_text, shapes):
        """Distribute watsonx text across spatial shapes"""
        if not watsonx_text.strip():
            return []
        
        # Simple distribution: split text into sentences and distribute
        sentences = [s.strip() for s in watsonx_text.split('.') if s.strip()]
        if not sentences:
            sentences = [watsonx_text]
        
        # Prioritize text-containing shapes
        text_shapes = [s for s in shapes if s.get('has_text_frame')]
        target_count = max(len(text_shapes), 1)
        
        # Distribute sentences across available slots
        text_parts = []
        sentences_per_part = max(1, len(sentences) // target_count)
        
        for i in range(0, len(sentences), sentences_per_part):
            part_sentences = sentences[i:i + sentences_per_part]
            text_part = '. '.join(part_sentences)
            if text_part and not text_part.endswith('.'):
                text_part += '.'
            text_parts.append(text_part)
        
        return text_parts



    def get_shape_type_name(self, shape_type):
        """Convert shape type enum to readable name"""
        if not PPTX_AVAILABLE:
            return f"Unknown_{shape_type}"
            
        shape_type_map = {
            MSO_SHAPE_TYPE.AUTO_SHAPE: "AutoShape",
            MSO_SHAPE_TYPE.CHART: "Chart", 
            MSO_SHAPE_TYPE.EMBEDDED_OLE_OBJECT: "EmbeddedOLEObject",   # Added for proper naming
            MSO_SHAPE_TYPE.GROUP: "Group",
            MSO_SHAPE_TYPE.LINE: "Line",
            MSO_SHAPE_TYPE.PICTURE: "Picture",
            MSO_SHAPE_TYPE.PLACEHOLDER: "Placeholder",
            MSO_SHAPE_TYPE.TABLE: "Table",
            MSO_SHAPE_TYPE.TEXT_BOX: "TextBox",
        }
        return shape_type_map.get(shape_type, f"Unknown_{shape_type}")



    def iter_textframed_shapes(self, shapes):
        """Generate shape objects that can contain text in document order"""
        for shape in shapes:
            # Recurse on group shapes
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                group_shape = shape
                for nested_shape in self.iter_textframed_shapes(group_shape.shapes):
                    yield nested_shape
                continue

            # Process leaf shapes with text frames
            if shape.has_text_frame:
                yield shape



    def extract_slide_text_structure(self, slide, slide_number):
        """Extract structured content from a single slide (from 08.ppt approach)"""
        slide_data = {
            'slide_number': slide_number,
            'title': '',
            'content': [],
            'notes': '',
            'layout_name': slide.slide_layout.name if hasattr(slide.slide_layout, 'name') else 'Unknown'
        }

        # Get all text-containing shapes
        textable_shapes = list(self.iter_textframed_shapes(slide.shapes))
        
        # Sort shapes by ROWS first (80-120% height tolerance), then LEFT-TO-RIGHT within rows
        print(f"      📖 Sorting {len(textable_shapes)} text shapes by ROWS first (80-120% height tolerance), then LEFT→RIGHT within rows")
        
        # Enhanced reading order with 80-120% height tolerance
        if textable_shapes:
            # FIXED: Use consistent 20% height tolerance for same row (80-120% range)
            avg_height = sum(shape.height for shape in textable_shapes) / len(textable_shapes)
            consistent_row_tolerance = max(avg_height * 0.2, 30)  # 20% tolerance for 80-120% height range
            
            def get_text_standard_reading_order_key(shape):
                # ROW-WISE READING: Use TOP edge for row determination (where content STARTS)
                shape_top = shape.top
                row_group = int(shape_top / consistent_row_tolerance) * consistent_row_tolerance
                return (
                    row_group,   # Primary: Row grouping based on TOP edge (where content starts)
                    shape.left   # Secondary: LEFT-TO-RIGHT within each row
                )
            
            ordered_shapes = sorted(textable_shapes, key=get_text_standard_reading_order_key)
            print(f"         ✅ Text shapes organized by ROWS first (80-120% height tolerance), then LEFT→RIGHT within rows")
        else:
            ordered_shapes = []

        # Process each shape in reading order
        for i, shape in enumerate(ordered_shapes):
            text_content = shape.text.strip()
            if not text_content:
                continue

            shape_info = {
                'reading_order': i + 1,
                'text': text_content,
                'position': {
                    'top': shape.top,
                    'left': shape.left,
                    'width': shape.width,
                    'height': shape.height
                },
                'shape_type': str(shape.shape_type).split('.')[-1] if hasattr(shape, 'shape_type') else 'UNKNOWN'
            }

            # Try to identify if this is likely a title (first shape, larger font, etc.)
            if i == 0 and not slide_data['title']:
                slide_data['title'] = text_content
                shape_info['is_title'] = True
            else:
                shape_info['is_title'] = False

            slide_data['content'].append(shape_info)

        # Extract slide notes if available
        try:
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                slide_data['notes'] = slide.notes_slide.notes_text_frame.text.strip()
        except:
            slide_data['notes'] = ''

        return slide_data



    def export_slides_as_text(self, text_structure):
        """Export each slide as a separate text file"""
        for slide in text_structure['slides']:
            filename = f"slide_{slide['slide_number']:02d}_{slide['title'][:30] if slide['title'] else 'untitled'}.txt"
            # Clean filename
            filename = "".join(c for c in filename if c.isalnum() or c in (' ', '-', '_', '.')).rstrip()
            file_path = self.text_structure_dir / filename

            with open(file_path, 'w', encoding='utf-8') as f:
                f.write(f"SLIDE {slide['slide_number']}\n")
                f.write(f"Layout: {slide['layout_name']}\n")
                f.write("=" * 50 + "\n\n")
                if slide['title']:
                    f.write(f"TITLE: {slide['title']}\n\n")
                f.write("CONTENT (Reading Order):\n")
                for item in slide['content']:
                    if not item['is_title']:
                        f.write(f"{item['reading_order']}. {item['text']}\n")
                if slide['notes']:
                    f.write(f"\nNOTES:\n{slide['notes']}\n")
        
        print(f"   📁 Individual slide text files saved to: {self.text_structure_dir}")



    def find_matching_watsonx_text(self, shape, watsonx_content):
        """Find matching watsonx text for a PowerPoint shape using content similarity"""
        if not watsonx_content or not hasattr(shape, 'text'):
            return None
        
        try:
            # Get PowerPoint shape text for comparison
            powerpoint_text = ""
            if hasattr(shape, 'text'):
                powerpoint_text = shape.text.strip()
            elif hasattr(shape, 'text_frame') and shape.text_frame:
                text_parts = []
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        text_parts.append(run.text)
                powerpoint_text = ''.join(text_parts).strip()
            
            if not powerpoint_text:
                return None
            
            # Search through watsonx content items for matches
            best_match = None
            best_score = 0.0
            extracted_text = None  # Initialize extracted text variable
            
            for content_item in watsonx_content:
                if not content_item.get('text'):
                    continue
                
                watsonx_text = content_item['text'].strip()
                extracted_text = None  # Reset for each content item
                
                # DEBUG: Show exact text comparison
                if "전륜" in powerpoint_text or "전륜" in watsonx_text:
                    print(f"      🔍 DETAILED COMPARISON:")
                    print(f"         PP: '{powerpoint_text}' (len={len(powerpoint_text)})")
                    print(f"         WX: '{watsonx_text}' (len={len(watsonx_text)})")
                    print(f"         PP repr: {repr(powerpoint_text)}")
                    print(f"         WX repr: {repr(watsonx_text)}")
                
                # SIMPLE WORD-BASED MATCHING: Count how many PP words are found in WX text
                pp_lower = powerpoint_text.lower().strip()
                wx_lower = watsonx_text.lower().strip()
                
                # Split into words and create word sets
                pp_words = pp_lower.split()
                wx_words_set = set(wx_lower.split())
                
                # Count matches: how many PP words are found in WX text
                matched_words = 0
                for pp_word in pp_words:
                    if pp_word in wx_words_set:
                        matched_words += 1
                
                # Calculate score: matched words / total PP words
                if pp_words:
                    score = matched_words / len(pp_words)
                    print(f"         → Word matching: {matched_words}/{len(pp_words)} PP words found in WX (score: {score:.3f})")
                else:
                    score = 0
                    print(f"         → No PP words to match (score: 0.000)")
                
                # Update best match if this score is better AND above 80% threshold
                if score > best_score and score >= 0.8:  # Must match at least 80% of PP words
                    best_score = score
                    best_match = watsonx_text  # Always use full watsonx text
                    print(f"      🔍 NEW BEST MATCH: PP='{powerpoint_text[:30]}...' → WX='{best_match[:50]}...' (score: {score:.2f})")
                elif score >= 0.8:
                    print(f"      🔍 GOOD MATCH: PP='{powerpoint_text[:30]}...' → WX='{watsonx_text[:50]}...' (score: {score:.2f})")
                else:
                    print(f"      🔍 POOR MATCH: Only {matched_words}/{len(pp_words)} words matched (score: {score:.3f})")
            
            if best_match:
                print(f"      🎯 FINAL MATCH: PowerPoint '{powerpoint_text[:30]}...' → watsonx text (score: {best_score:.2f})")
                return best_match
            else:
                print(f"      ❌ NO MATCH: PowerPoint '{powerpoint_text[:30]}...' → no watsonx match found")
            
            return None
            
        except Exception as e:
            print(f"      ⚠️  Error matching watsonx text: {e}")
            return None


